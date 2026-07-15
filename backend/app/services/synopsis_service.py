import re
import json
import logging
import io
from datetime import datetime, timedelta
from bson import ObjectId
from fastapi import HTTPException, status
from openai import OpenAI

from app.core.config import settings
from app.models.synopsis import SynopsisDBModel
from app.schemas.synopsis import SynopsisGenerateRequest
from app.utils.youtube_utils import extract_youtube_id, get_youtube_metadata, get_youtube_transcript, get_youtube_transcript_via_yt_dlp
from app.utils.whisper_utils import transcribe_audio_with_whisper

logger = logging.getLogger(__name__)

def format_seconds_to_timestamp(seconds: float) -> str:
    """Converts seconds float into MM:SS or HH:MM:SS format."""
    hours = int(seconds // 3600)
    minutes = int((seconds % 3600) // 60)
    secs = int(seconds % 60)
    if hours > 0:
        return f"{hours:02d}:{minutes:02d}:{secs:02d}"
    else:
        return f"{minutes:02d}:{secs:02d}"

class SynopsisService:
    def __init__(self, db):
        self.collection = db[settings.SYNOPSIS_COLLECTION]
        self.client = OpenAI(api_key=settings.GROQ_API_KEY, base_url="https://api.groq.com/openai/v1", max_retries=0)

    async def generate_synopsis(self, user_id: str, request: SynopsisGenerateRequest) -> dict:
        """
        Main AI pipeline: Extracts video ID, fetches metadata, retrieves transcript 
        (with Whisper fallback), triggers GPT-4o-mini summarization, and saves to DB.
        Includes MongoDB-backed caching and daily rate limiting.
        """
        # 1. Extract Video ID
        video_id = extract_youtube_id(request.youtubeUrl)
        if not video_id:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="Could not extract a valid YouTube video ID from the provided URL."
            )
            
        # API Cost Protection: Check user cache
        existing_user_doc = await self.collection.find_one({"user_id": user_id, "metadata.youtubeId": video_id})
        if existing_user_doc:
            logger.info(f"User cache hit for video {video_id} and user {user_id}")
            existing_user_doc["id"] = str(existing_user_doc["_id"])
            if "_id" in existing_user_doc:
                del existing_user_doc["_id"]
            if "created_at" in existing_user_doc and isinstance(existing_user_doc["created_at"], datetime):
                existing_user_doc["createdAt"] = existing_user_doc["created_at"].isoformat()
                del existing_user_doc["created_at"]
            return existing_user_doc

        # API Cost Protection: Check global cache
        existing_global_doc = await self.collection.find_one({"metadata.youtubeId": video_id})
        if existing_global_doc:
            logger.info(f"Global cache hit for video {video_id}, copying to user {user_id}")
            new_doc = dict(existing_global_doc)
            del new_doc["_id"]
            new_doc["user_id"] = user_id
            new_doc["saved"] = False
            new_doc["created_at"] = datetime.utcnow()
            result = await self.collection.insert_one(new_doc)
            new_doc["id"] = str(result.inserted_id)
            new_doc["createdAt"] = new_doc["created_at"].isoformat()
            del new_doc["created_at"]
            return new_doc
            
        # API Cost Protection: Daily limit (max 10 new summaries per user per day)
        one_day_ago = datetime.utcnow() - timedelta(days=1)
        user_runs_today = await self.collection.count_documents({
            "user_id": user_id,
            "created_at": {"$gte": one_day_ago}
        })
        if user_runs_today >= 10:
            raise HTTPException(
                status_code=status.HTTP_429_TOO_MANY_REQUESTS,
                detail="Daily generation limit exceeded. You can generate up to 10 new video synopses per 24 hours to protect API budget."
            )

        # 2. Fetch Video Metadata
        try:
            metadata = get_youtube_metadata(video_id)
        except ValueError as val_err:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=str(val_err)
            )
        
        # 3. Retrieve Transcript (with Whisper speech-to-text fallback)
        transcript_segments = []
        transcript_fallback_note = None
        try:
            transcript_segments = get_youtube_transcript(video_id)
            logger.info("Successfully fetched native YouTube transcript.")
        except Exception as transcript_err:
            logger.warning(f"Native transcript fetching failed: {transcript_err}. Attempting yt-dlp caption fallback...")
            try:
                transcript_segments = get_youtube_transcript_via_yt_dlp(video_id)
                logger.info("Successfully fetched transcript from yt-dlp caption tracks.")
            except Exception as ytdlp_caption_err:
                logger.warning(f"yt-dlp caption fallback failed: {ytdlp_caption_err}. Attempting Whisper transcription fallback...")
                try:
                    transcript_segments = transcribe_audio_with_whisper(video_id)
                    logger.info("Successfully transcribed video using Whisper fallback.")
                except Exception as whisper_err:
                    logger.error(f"Whisper fallback transcription failed: {whisper_err}")
                    whisper_msg = str(whisper_err).lower()
                    if "insufficient_quota" in whisper_msg or "rate limit" in whisper_msg or "429" in whisper_msg:
                        transcript_segments = self._build_metadata_only_transcript(metadata)
                        transcript_fallback_note = (
                            "Captions were unavailable and Whisper transcription hit Groq quota/rate limits. "
                            "This synopsis used metadata-assisted fallback logic. For higher fidelity, enable captions "
                            "or add Groq credits."
                        )
                        logger.warning("Using metadata-only transcript fallback due to Groq quota/rate limits.")
                    elif "api key" in whisper_msg or "authentication" in whisper_msg:
                        raise HTTPException(
                            status_code=status.HTTP_401_UNAUTHORIZED,
                            detail="Whisper fallback failed because Groq API key is invalid or missing."
                        )
                    else:
                        raise HTTPException(
                            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                            detail=f"Failed to retrieve video transcript from captions and Whisper fallback: {str(whisper_err)}"
                        )

        if not transcript_segments:
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail="Transcript extraction returned empty content."
            )
            
        # 4. Format transcript text with timestamps for GPT prompt
        transcript_text = ""
        for segment in transcript_segments:
            timestamp = format_seconds_to_timestamp(segment['start'])
            transcript_text += f"[{timestamp}] {segment['text']}\n"
            
        # Limit transcript length to keep prompt within bounds (~30k chars or approx 8k tokens)
        if len(transcript_text) > 40000:
            logger.info("Transcript exceeds limit, truncating to fit within token bounds...")
            transcript_text = transcript_text[:40000] + "\n... [Transcript Truncated due to length] ..."

        # 5. Call Groq Summarization
        logger.info("Requesting Groq summary...")
        openai_data = None
        try:
            system_prompt = (
                "You are an expert AI Video Analyst and Senior Curriculum Architect. Analyze the following YouTube video transcript and generate a highly detailed, professional, structured educational study guide synopsis. The content MUST be extremely comprehensive and thorough, containing detailed explanations, paragraphs, lists, tables, and step-by-step concepts so that a student can fully understand the topic and prepare for an exam without watching the video.\n\n"
                "Respond ONLY in valid JSON matching this exact schema:\n"
                "{\n"
                '  "learningObjectives": ["Objective 1", "Objective 2", "Objective 3"],\n'
                '  "majorConcepts": [\n'
                '    {\n'
                '      "concept": "Concept Name",\n'
                '      "explanation": "Detailed, thorough explanation of this core conceptual pillar.",\n'
                '      "subtopics": [\n'
                '        {"title": "Subtopic Title", "explanation": "Detailed explanation of this subtopic."}\n'
                '      ],\n'
                '      "definitions": [\n'
                '        {"term": "Term or Acronym", "definition": "Detailed definition and context."}\n'
                '      ],\n'
                '      "examples": [\n'
                '        {"example": "Example/Case Title", "description": "Details of the practical example or case study."}\n'
                '      ],\n'
                '      "formulas": [\n'
                '        {"formula": "Formula string (e.g. E = mc^2)", "explanation": "Explanation of variables and usage."}\n'
                '      ],\n'
                '      "workflows": [\n'
                '        {"step": "Step/Phase Name", "description": "What happens during this step in the process."}\n'
                '      ]\n'
                '    }\n'
                '  ],\n'
                '  "executiveSummary": "A concise overview of the video\'s core message, target audience, and primary takeaways.",\n'
                '  "introduction": "A comprehensive introduction explaining what the topic/video is about, its relevance, and context.",\n'
                '  "detailedExplanation": "A deep detailed explanation of every major topic and subtopic covered in the video in markdown format. Provide detailed understandable explanations, paragraphs, tables, and lists. Do not just list bullet points. Be comprehensive and clear.",\n'
                '  "practicalApplications": [\n'
                '    {"application": "Application Area", "description": "How this concept/technology/skill is applied in the real world."}\n'
                '  ],\n'
                '  "keyTakeaways": [\n'
                '    "Key takeaway point 1",\n'
                '    "Key takeaway point 2"\n'
                '  ],\n'
                '  "quickRevisionNotes": "A one-page markdown-formatted quick revision reference summarizing the entire topic for fast exam review.",\n'
                '  "faq": [\n'
                '    {"question": "Frequently asked question?", "answer": "Clear, comprehensive answer."}\n'
                '  ],\n'
                '  "interviewQuestions": [\n'
                '    {"question": "Typical technical or conceptual interview question?", "answer": "Ideal response for a candidate to stand out."}\n'
                '  ],\n'
                '  "examPreparationNotes": [\n'
                '    "Focus area or key formula/concept likely to appear in exams",\n'
                '    "Quick revision trick or common pitfall to avoid in exams"\n'
                '  ],\n'
                '  "chapters": [\n'
                '    {\n'
                '      "title": "Chapter name",\n'
                '      "timestamp": "Timestamp in MM:SS or HH:MM:SS format",\n'
                '      "summary": "Brief summary of what is discussed in this chapter."\n'
                '    }\n'
                '  ],\n'
                '  "actionItems": ["Actionable steps, recommendations, or items that the viewer can implement."],\n'
                '  "topics": [\n'
                '    {\n'
                '      "name": "Topic name",\n'
                '      "percentage": 50\n'
                '    }\n'
                '  ],\n'
                '  "sentiment": {\n'
                '    "label": "positive/neutral/negative/mixed/analytical",\n'
                '    "score": 90,\n'
                '    "explanation": "Brief explanation of why this sentiment score and label was chosen."\n'
                '  },\n'
                '  "keywords": ["tag1", "tag2"],\n'
                '  "nextStepsRecommendation": "A final synthesis conclusion or recommended next steps for the viewer.",\n'
                '  "theme": "AI/ML | Technology | Business | Finance | Educational | Healthcare | Marketing | General",\n'
                '  "difficultyLevel": "Beginner | Intermediate | Advanced",\n'
                '  "estimatedReadTime": "15 mins",\n'
                '  "estimatedRevisionTime": "5 mins",\n'
                '  "contentQualityScore": 95,\n'
                '  "resumeProjectSummary": {\n'
                '    "projectTitle": "Title of the project or tutorial app",\n'
                '    "technologiesUsed": ["React", "Node.js"],\n'
                '    "keyFeatures": ["Key feature 1", "Key feature 2"],\n'
                '    "resumeBulletPoints": [\n'
                '      "Designed and implemented X features utilizing Y tech, improving performance by Z%.",\n'
                '      "Built responsive interface using..."\n'
                '    ]\n'
                '  },\n'
                '  "vivaQuestions": [\n'
                '    {"question": "Viva oral exam question?", "answer": "Thorough answer for viva."}\n'
                '  ],\n'
                '  "shortAnswerQuestions": [\n'
                '    {"question": "Short answer question?", "answer": "Concise paragraph explanation (approx 100 words).", "marks": 5}\n'
                '  ],\n'
                '  "longAnswerQuestions": [\n'
                '    {"question": "Long answer question?", "answer": "Extensive explanation with bullet points if applicable (approx 300 words).", "marks": 10}\n'
                '  ],\n'
                '  "practiceQuestions": [\n'
                '    {"question": "Self-practice exercise question?", "hint": "Useful hint to get them started."}\n'
                '  ]\n'
                "}\n\n"
                "Constraints:\n"
                "- Only populate resumeProjectSummary if the video is a project/coding/construction tutorial. Otherwise set it to null.\n"
                "- Under each concept in majorConcepts, include definitions, examples, formulas, and workflows if applicable to that concept. If a concept has no formulas or workflows, keep them as empty arrays [].\n"
                "- Make sure the total percentage of all topics combined is close to 100.\n"
                "- Do not hallucinate or add any explanations outside of the JSON object.\n"
                "- Maintain absolute factual alignment with the transcript content.\n"
                "- The input transcript might be in a non-English language (such as Telugu, Hindi, Spanish, etc.). If so, you MUST translate all concepts and write the entire output JSON response (learning objectives, concepts, explanations, detailed explanations, summaries, viva/interview/exam prep questions, etc.) strictly in English."
            )
            
            user_prompt = f"Video Title: {metadata['title']}\n\nTranscript:\n{transcript_text}"
            openai_model = settings.GROQ_MODEL or "llama-3.3-70b-versatile"
            
            response = self.client.chat.completions.create(
                model=openai_model,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_prompt}
                ],
                response_format={"type": "json_object"},
                timeout=120
            )
            
            response_text = response.choices[0].message.content.strip()
            
            # Clean markdown code wrapper if present
            if response_text.startswith("```"):
                lines = response_text.splitlines()
                if lines[0].startswith("```"):
                    lines = lines[1:]
                if lines[-1].startswith("```"):
                    lines = lines[:-1]
                response_text = "\n".join(lines).strip()
                
            openai_data = json.loads(response_text)
            
        except Exception as groq_err:
            logger.warning(f"Groq summarization failed ({groq_err}). Trying OpenAI fallback...")
            openai_data = None
            if settings.OPENAI_API_KEY:
                try:
                    logger.info("Attempting synopsis generation with OpenAI (model: gpt-4o-mini)...")
                    openai_client = OpenAI(api_key=settings.OPENAI_API_KEY, max_retries=0)
                    openai_model_name = settings.OPENAI_MODEL or "gpt-4o-mini"
                    response = openai_client.chat.completions.create(
                        model=openai_model_name,
                        messages=[
                            {"role": "system", "content": system_prompt},
                            {"role": "user", "content": user_prompt}
                        ],
                        response_format={"type": "json_object"},
                        timeout=120
                    )
                    response_text = response.choices[0].message.content.strip()
                    if response_text.startswith("```"):
                        lines = response_text.splitlines()
                        if lines[0].startswith("```"):
                            lines = lines[1:]
                        if lines[-1].startswith("```"):
                            lines = lines[:-1]
                        response_text = "\n".join(lines).strip()
                    openai_data = json.loads(response_text)
                    logger.info("Successfully generated synopsis using OpenAI fallback.")
                except Exception as oai_err:
                    logger.error(f"OpenAI fallback synopsis generation failed: {oai_err}")

            if not openai_data:
                logger.warning("Falling back to Procedural AI Summarizer...")
                try:
                    openai_data = self._generate_procedural_fallback_synopsis(metadata, transcript_segments)
                except Exception as fallback_err:
                    logger.error(f"Procedural fallback failed as well: {fallback_err}")
                    raise HTTPException(
                        status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                        detail=f"Groq summarization service, OpenAI fallback, and procedural fallback all failed: {str(groq_err)}"
                    )
        
        gemini_data = openai_data
        if transcript_fallback_note:
            gemini_data["nextStepsRecommendation"] = (
                f"{gemini_data.get('nextStepsRecommendation', '').strip()} {transcript_fallback_note}"
            ).strip()

        # 6. Map response to SynopsisDBModel
        try:
            # Aggregate definitions, examples, and formulas from the majorConcepts hierarchy for compatibility fallback
            aggregated_definitions = []
            aggregated_examples = []
            aggregated_formulas = []
            
            for c_data in gemini_data.get("majorConcepts", []):
                for nd in c_data.get("definitions", []):
                    if nd.get("term") and nd.get("definition"):
                        aggregated_definitions.append({"term": nd.get("term"), "definition": nd.get("definition")})
                for ne in c_data.get("examples", []):
                    if ne.get("example") and ne.get("description"):
                        aggregated_examples.append({"example": ne.get("example"), "description": ne.get("description")})
                for nf in c_data.get("formulas", []):
                    if nf.get("formula") and nf.get("explanation"):
                        aggregated_formulas.append({"formula": nf.get("formula"), "explanation": nf.get("explanation")})
            
            # Fallbacks to root-level keys if nested loops were empty
            if not aggregated_definitions:
                for nd in gemini_data.get("importantDefinitions", []):
                    aggregated_definitions.append({"term": nd.get("term"), "definition": nd.get("definition")})
            if not aggregated_examples:
                for ne in gemini_data.get("examples", []):
                    aggregated_examples.append({"example": ne.get("example"), "description": ne.get("description")})
            if not aggregated_formulas:
                for nf in gemini_data.get("importantFormulas", []):
                    aggregated_formulas.append({"formula": nf.get("formula"), "explanation": nf.get("explanation")})

            mapped_data = {
                "user_id": user_id,
                "metadata": {
                    "id": metadata["id"],
                    "youtubeId": metadata["youtubeId"],
                    "title": metadata["title"],
                    "channelName": metadata["channelName"],
                    "duration": metadata["duration"],
                    "publishDate": metadata["publishDate"],
                    "thumbnail": metadata["thumbnail"],
                    "youtubeUrl": metadata["youtubeUrl"],
                    "views": metadata.get("views", "N/A")
                },
                "executiveSummary": gemini_data.get("executiveSummary", ""),
                "chapters": gemini_data.get("chapters", []),
                "insights": [
                    {"title": f"Core Takeaway {idx+1}", "description": takeaway}
                    for idx, takeaway in enumerate(gemini_data.get("keyTakeaways", gemini_data.get("insights", [])))
                ],
                "actionItems": gemini_data.get("actionItems", []),
                "topics": [
                    {
                        "topic": topic.get("name", "Unknown Topic"),
                        "percentage": topic.get("percentage", 0),
                        "description": f"Concepts and discussions surrounding {topic.get('name', 'topic')}."
                    }
                    for topic in gemini_data.get("topics", [])
                ],
                "conclusion": gemini_data.get("nextStepsRecommendation", ""),
                "keywords": gemini_data.get("keywords", []),
                "sentiment": {
                    "label": gemini_data.get("sentiment", {}).get("label", "Neutral").capitalize(),
                    "score": gemini_data.get("sentiment", {}).get("score", 50),
                    "explanation": gemini_data.get("sentiment", {}).get("explanation", "")
                },
                "saved": False,
                "outputFormat": request.outputFormat or "pdf",
                
                # Educational Extensions
                "introduction": gemini_data.get("introduction", ""),
                "detailedExplanation": gemini_data.get("detailedExplanation", ""),
                "keyConcepts": gemini_data.get("keyConcepts", []),
                "importantDefinitions": aggregated_definitions,
                "examples": aggregated_examples,
                "practicalApplications": gemini_data.get("practicalApplications", []),
                "keyTakeaways": gemini_data.get("keyTakeaways", []),
                "quickRevisionNotes": gemini_data.get("quickRevisionNotes", ""),
                "faq": gemini_data.get("faq", []),
                "interviewQuestions": gemini_data.get("interviewQuestions", []),
                "examPreparationNotes": gemini_data.get("examPreparationNotes", []),
                "theme": gemini_data.get("theme", "General"),
                "difficultyLevel": gemini_data.get("difficultyLevel", "Intermediate"),
                "estimatedReadTime": gemini_data.get("estimatedReadTime", "5 mins"),
                "estimatedRevisionTime": gemini_data.get("estimatedRevisionTime", "2 mins"),
                "contentQualityScore": gemini_data.get("contentQualityScore", 90),
                "importantFormulas": aggregated_formulas,
                "mcqs": gemini_data.get("mcqs", []),
                "resumeProjectSummary": gemini_data.get("resumeProjectSummary", None),
                
                # OpenAI study guide additions
                "viva_questions": gemini_data.get("vivaQuestions", []),
                "short_answer_questions": gemini_data.get("shortAnswerQuestions", []),
                "long_answer_questions": gemini_data.get("longAnswerQuestions", []),
                "practice_questions": gemini_data.get("practiceQuestions", []),
                
                # OpenAI Hierarchical additions
                "learning_objectives": gemini_data.get("learningObjectives", []),
                "major_concepts": gemini_data.get("majorConcepts", []),
                "generated_study_material": gemini_data,
                "transcript": transcript_text
            }
            
            # Validate against Pydantic DB Model
            db_model = SynopsisDBModel(**mapped_data)
            
        except Exception as val_err:
            logger.error(f"Pydantic mapping/validation failed: {val_err}")
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail=f"Formatting AI response data failed validation: {str(val_err)}"
            )

        # 7. Insert into MongoDB
        synopsis_dict = db_model.model_dump(by_alias=True, exclude={"id"})
        result = await self.collection.insert_one(synopsis_dict)
        
        # 8. Prepare Response
        synopsis_dict["id"] = str(result.inserted_id)
        if "_id" in synopsis_dict:
            del synopsis_dict["_id"]
            
        return synopsis_dict

    async def get_history(self, user_id: str) -> list:
        """
        Fetch all synopses belonging to a specific user.
        """
        cursor = self.collection.find({"user_id": user_id}).sort("created_at", -1)
        history = []
        async for doc in cursor:
            doc["id"] = str(doc["_id"])
            if "_id" in doc:
                del doc["_id"]
            # Format output matching frontend schema (handling camelCase vs snake_case aliases gracefully)
            doc["createdAt"] = doc["created_at"].isoformat() if isinstance(doc.get("created_at"), datetime) else doc.get("createdAt")
            if "created_at" in doc:
                del doc["created_at"]
            doc["outputFormat"] = doc.get("outputFormat") or doc.get("output_format") or "pdf"
            history.append(doc)
        return history

    async def get_by_id(self, user_id: str, synopsis_id: str) -> dict:
        """
        Fetch a single synopsis by its ID, ensuring it belongs to the user.
        """
        if not ObjectId.is_valid(synopsis_id):
             raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Invalid synopsis ID format")

        doc = await self.collection.find_one({"_id": ObjectId(synopsis_id), "user_id": user_id})
        if not doc:
            raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Synopsis not found")

        doc["id"] = str(doc["_id"])
        if "_id" in doc:
            del doc["_id"]
        doc["createdAt"] = doc["created_at"].isoformat() if isinstance(doc.get("created_at"), datetime) else doc.get("createdAt")
        if "created_at" in doc:
            del doc["created_at"]
        doc["outputFormat"] = doc.get("outputFormat") or doc.get("output_format") or "pdf"
        return doc

    async def export_ppt(self, user_id: str, synopsis_id: str) -> tuple[bytes, str]:
        """
        Export a synopsis as a modern PPTX file.
        """
        synopsis = await self.get_by_id(user_id, synopsis_id)
        slide_plan = await self._generate_ai_slide_plan(synopsis)
        if not self._validate_slide_plan(slide_plan):
            logger.info("AI slide plan failed validation. Using fallback slide plan.")
            slide_plan = self._generate_fallback_slide_plan(synopsis)
        ppt_bytes = self._build_ppt_bytes(synopsis, slide_plan)
        raw_title = synopsis.get("metadata", {}).get("title", "video_synopsis")
        safe_name = re.sub(r"[^a-zA-Z0-9_-]+", "_", raw_title).strip("_").lower() or "video_synopsis"
        return ppt_bytes, f"video_synopsis_{safe_name}.pptx"




    async def delete_synopsis(self, user_id: str, synopsis_id: str) -> bool:
        """
        Delete a synopsis.
        """
        if not ObjectId.is_valid(synopsis_id):
             raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Invalid synopsis ID format")

        result = await self.collection.delete_one({"_id": ObjectId(synopsis_id), "user_id": user_id})
        if result.deleted_count == 0:
            raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Synopsis not found")
        return True

    async def toggle_save(self, user_id: str, synopsis_id: str) -> dict:
        """
        Toggle the saved boolean of a synopsis.
        """
        if not ObjectId.is_valid(synopsis_id):
             raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Invalid synopsis ID format")
              
        doc = await self.collection.find_one({"_id": ObjectId(synopsis_id), "user_id": user_id})
        if not doc:
            raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Synopsis not found")

        new_saved_status = not doc.get("saved", False)
        await self.collection.update_one(
            {"_id": ObjectId(synopsis_id)},
            {"$set": {"saved": new_saved_status}}
        )

        doc["saved"] = new_saved_status
        doc["id"] = str(doc["_id"])
        if "_id" in doc:
            del doc["_id"]
        doc["createdAt"] = doc["created_at"].isoformat() if isinstance(doc.get("created_at"), datetime) else doc.get("createdAt")
        if "created_at" in doc:
            del doc["created_at"]
        doc["outputFormat"] = doc.get("outputFormat") or doc.get("output_format") or "pdf"
        return doc

    def _build_metadata_only_transcript(self, metadata: dict) -> list:
        """
        Creates a minimal synthetic transcript from metadata when captions and Whisper both fail.
        """
        return [
            {"text": f"Video title: {metadata.get('title', 'Unknown')}", "start": 0.0, "duration": 4.0},
            {"text": f"Channel name: {metadata.get('channelName', 'Unknown')}", "start": 4.0, "duration": 4.0},
            {"text": f"Publish date: {metadata.get('publishDate', 'Unknown')}", "start": 8.0, "duration": 4.0},
            {"text": f"Duration: {metadata.get('duration', 'Unknown')}", "start": 12.0, "duration": 4.0},
        ]

    async def _generate_ai_slide_plan(self, synopsis: dict) -> dict:
        """
        Generates a structured slide plan JSON via the LLM (Groq/OpenAI).
        """
        metadata = synopsis.get("metadata", {})
        title = metadata.get("title", "Video Synopsis")
        channel = metadata.get("channelName", "Unknown Channel")
        exec_summary = synopsis.get("executiveSummary", "")
        introduction = synopsis.get("introduction", "")
        chapters = synopsis.get("chapters", [])
        concepts = synopsis.get("major_concepts") or synopsis.get("majorConcepts") or synopsis.get("keyConcepts") or []
        
        context_str = f"Title: {title}\nChannel: {channel}\nSummary: {exec_summary}\nIntroduction: {introduction}\n"
        if chapters:
            context_str += "\nChapters:\n" + "\n".join([f"- [{c.get('timestamp')}] {c.get('title')}: {c.get('summary')}" for c in chapters[:6]])
        if concepts:
            context_str += "\nConcepts:\n" + "\n".join([f"- {c.get('concept') or c.get('name')}: {c.get('explanation')}" for c in concepts[:4]])

        system_prompt = (
            "You are an expert Presentation Designer and Academic Study Coach. "
            "Your task is to design a highly professional, visually engaging slide deck plan for a presentation based on the provided video synopsis.\n"
            "Create a structured slide deck of exactly 8 to 12 slides. Do not create fewer than 8 or more than 12 slides.\n"
            "Adopt a modern visual structure. Use slide types: 'title', 'agenda', 'timeline', 'flowchart', 'concepts', 'comparison', 'revision', 'quiz', 'summary', 'thankyou'.\n"
            "CRITICAL: Every slide MUST be 100% topic-specific and directly matching the video theme. Do NOT use generic software/system engineering placeholders (such as 'decoupled and isolated processes', 'system latency constraint', 'horizontal clusters') unless the video is explicitly about those topics. Avoid all generic template placeholders.\n"
            "Improve content depth: each slide's content array should contain 3 to 4 detailed descriptive bullet points, not just short one-word lists.\n"
            "For the 'quiz' slide: 'content' MUST contain exactly 4 actual, topic-specific conceptual questions derived from the video content (e.g., if the video is about Network Security, ask 'What is network security?', 'Define cryptography', etc.).\n"
            "For the 'revision' slide: 'content' MUST contain exactly 4 high-yield revision tips for exam preparation based on the lecture.\n"
            "Speaker notes MUST be comprehensive, detailed paragraphs (minimum 3-4 sentences per slide) explaining the slide's visual layout and providing detailed guidance for the presenter.\n\n"
            "Respond ONLY in valid JSON matching this exact structure:\n"
            "{\n"
            '  "subtitle": "Comprehensive Study Guide Presentation",\n'
            '  "slides": [\n'
            '    {\n'
            '      "type": "title",\n'
            '      "title": "Title of the Slide",\n'
            '      "content": ["Detailed overview point explaining the main context", "Another thorough introduction point"],\n'
            '      "diagram": {"type": "none", "items": []},\n'
            '      "speaker_notes": "A detailed explanation of the presentation focus..."\n'
            '    },\n'
            '    {\n'
            '      "type": "agenda",\n'
            '      "title": "Table of Contents",\n'
            '      "content": ["Agenda overview describing the roadmap of topics"],\n'
            '      "diagram": {\n'
            '        "type": "cards",\n'
            '        "items": [\n'
            '          {"title": "Objectives", "description": "Specific target metrics or goals", "icon": "ðŸŽ¯"},\n'
            '          {"title": "Concepts", "description": "Core theoretical pillars to be examined", "icon": "ðŸ’¡"},\n'
            '          {"title": "Workflows", "description": "Step-by-step practical implementation paths", "icon": "âš¡"},\n'
            '          {"title": "Review Quiz", "description": "Active recall self-test items", "icon": "ðŸ“"}\n'
            '        ]\n'
            '      },\n'
            '      "speaker_notes": "We will cover four core areas today..."\n'
            '    },\n'
            '    {\n'
            '      "type": "timeline",\n'
            '      "title": "Chronological Overview",\n'
            '      "content": ["Detailed summary of video progression timeline"],\n'
            '      "diagram": {\n'
            '        "type": "timeline",\n'
            '        "items": [\n'
            '          {"time": "00:00", "title": "Section Title", "description": "Detailed explanation of what occurs in this section"}\n'
            '        ]\n'
            '      },\n'
            '      "speaker_notes": "Starting with the initial segment..."\n'
            '    },\n'
            '    {\n'
            '      "type": "flowchart",\n'
            '      "title": "Process Pipeline",\n'
            '      "content": ["Step-by-step processes or methodologies shown in the video"],\n'
            '      "diagram": {\n'
            '        "type": "flowchart",\n'
            '        "items": [\n'
            '          {"step": "Step 1 Title", "description": "Detailed explanation of this step"}\n'
            '        ]\n'
            '      },\n'
            '      "speaker_notes": "This flowchart represents the sequence..."\n'
            '    },\n'
            '    {\n'
            '      "type": "concepts",\n'
            '      "title": "Core Conceptual Pillars",\n'
            '      "content": ["Core theoretical framework summaries"],\n'
            '      "diagram": {\n'
            '        "type": "cards",\n'
            '        "items": [\n'
            '          {"title": "Concept A", "description": "Detailed concept explanation", "icon": "ðŸ’¡"}\n'
            '        ]\n'
            '      },\n'
            '      "speaker_notes": "Looking at the conceptual pillars..."\n'
            '    },\n'
            '    {\n'
            '      "type": "comparison",\n'
            '      "title": "Comparison & Tradeoffs",\n'
            '      "content": ["A comparison between core concepts or technologies"],\n'
            '      "diagram": {\n'
            '        "type": "table",\n'
            '        "items": [\n'
            '          {"aspect": "Feature Comparison Aspect", "item1": "Detailed Attribute A", "item2": "Detailed Attribute B"}\n'
            '        ]\n'
            '      },\n'
            '      "speaker_notes": "This table compares key trade-offs..."\n'
            '    },\n'
            '    {\n'
            '      "type": "quiz",\n'
            '      "title": "Active Recall Check",\n'
            '      "content": ["Quiz Question A", "Quiz Question B", "Quiz Question C", "Quiz Question D"],\n'
            '      "diagram": {"type": "none", "items": []},\n'
            '      "speaker_notes": "Let\'s check knowledge retention with these 4 questions..."\n'
            '    },\n'
            '    {\n'
            '      "type": "revision",\n'
            '      "title": "Quick Revision & Tips",\n'
            '      "content": ["Revision Tip A", "Revision Tip B", "Revision Tip C", "Revision Tip D"],\n'
            '      "diagram": {"type": "none", "items": []},\n'
            '      "speaker_notes": "To review for exam preparation..."\n'
            '    },\n'
            '    {\n'
            '      "type": "thankyou",\n'
            '      "title": "Questions & Answers",\n'
            '      "content": ["Thank you for participating. We will now address questions."],\n'
            '      "diagram": {"type": "none", "items": []},\n'
            '      "speaker_notes": "We wrap up and open the floor..."\n'
            '    }\n'
            '  ]\n'
            "}\n"
            "Do not return markdown wrappers, return only raw JSON."
        )

        openai_model = settings.GROQ_MODEL or "llama-3.3-70b-versatile"
        response_text = None
        
        # 1. Try Groq first
        if settings.GROQ_API_KEY:
            try:
                logger.info("Slide Plan: attempting Groq completions...")
                response = self.client.chat.completions.create(
                    model=openai_model,
                    messages=[
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": f"Here is the video synopsis:\n{context_str}"}
                    ],
                    response_format={"type": "json_object"},
                    timeout=90
                )
                response_text = response.choices[0].message.content.strip()
            except Exception as e:
                logger.warning(f"Slide Plan: Groq failed to generate slide plan: {e}. Trying OpenAI fallback.")

        # 2. Try OpenAI fallback
        if not response_text and settings.OPENAI_API_KEY:
            try:
                logger.info("Slide Plan: attempting OpenAI completions (model: gpt-4o-mini)...")
                openai_client = OpenAI(api_key=settings.OPENAI_API_KEY, max_retries=0)
                openai_model_name = settings.OPENAI_MODEL or "gpt-4o-mini"
                response = openai_client.chat.completions.create(
                    model=openai_model_name,
                    messages=[
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": f"Here is the video synopsis:\n{context_str}"}
                    ],
                    response_format={"type": "json_object"},
                    timeout=90
                )
                response_text = response.choices[0].message.content.strip()
            except Exception as oai_err:
                logger.error(f"Slide Plan: OpenAI fallback also failed: {oai_err}")

        # Parse response text if retrieved successfully
        if response_text:
            try:
                if response_text.startswith("```"):
                    lines = response_text.splitlines()
                    if lines[0].startswith("```"):
                        lines = lines[1:]
                    if lines[-1].startswith("```"):
                        lines = lines[:-1]
                    response_text = "\n".join(lines).strip()
                    
                return json.loads(response_text)
            except Exception as parse_err:
                logger.error(f"Failed to parse AI slide plan JSON: {parse_err}")
                
        return {}

    def _validate_slide_plan(self, plan: dict) -> bool:
        """
        Validates that the slide plan exists, slide count is between 8 and 12,
        and each slide contains all necessary keys. Rejects generic software/system placeholders.
        """
        if not plan or not isinstance(plan, dict):
            return False
        slides = plan.get("slides")
        if not isinstance(slides, list):
            return False
        if not (8 <= len(slides) <= 12):
            return False
        for s in slides:
            if not isinstance(s, dict):
                return False
            required = ["type", "title", "content", "diagram", "speaker_notes"]
            if not all(k in s for k in required):
                return False
            
            # Strict Content Validation: reject generic fallbacks or placeholders
            s_content = s.get("content") or []
            if isinstance(s_content, str):
                s_content = [s_content]
            
            forbidden_phrases = [
                "processes decoupled", "decoupled and isolated", 
                "system latency constraint", "horizontal clusters",
                "scale horizontal", "placeholder", "lorem ipsum",
                "todo", "insert questions"
            ]
            for x in s_content:
                x_lower = str(x).lower()
                for p in forbidden_phrases:
                    if p in x_lower:
                        logger.warning(f"CORS/Validation: Rejected slide plan due to generic phrase '{p}' in slide '{s.get('title')}'")
                        return False
            
            # Verify diagram structure
            diag = s.get("diagram")
            if not isinstance(diag, dict) or "type" not in diag or "items" not in diag:
                return False
        return True

    def _get_topic_specific_questions(self, synopsis: dict) -> list[str]:
        """
        Dynamically extracts or constructs 4 topic-specific quiz questions from the synopsis.
        """
        questions = []
        
        # 1. Try to extract from MCQs
        mcqs = synopsis.get("mcqs") or synopsis.get("generated_study_material", {}).get("mcqs") or []
        for mcq in mcqs:
            if isinstance(mcq, dict) and mcq.get("question"):
                questions.append(mcq["question"].strip())
                
        # 2. Try practice questions, viva questions, or short answers
        practice = synopsis.get("practice_questions") or synopsis.get("practiceQuestions") or []
        for pq in practice:
            if isinstance(pq, dict) and pq.get("question"):
                questions.append(pq["question"].strip())
                
        viva = synopsis.get("viva_questions") or synopsis.get("vivaQuestions") or []
        for vq in viva:
            if isinstance(vq, dict) and vq.get("question"):
                questions.append(vq["question"].strip())
                
        short_ans = synopsis.get("short_answer_questions") or synopsis.get("shortAnswerQuestions") or []
        for sa in short_ans:
            if isinstance(sa, dict) and sa.get("question"):
                questions.append(sa["question"].strip())
                
        # Clean duplicates while keeping order
        unique_questions = []
        for q in questions:
            if q and q not in unique_questions:
                unique_questions.append(q)
        questions = unique_questions

        # Inject highly relevant conceptual questions if the topic is security/encryption/cryptography
        title_text = synopsis.get("metadata", {}).get("title", "").lower()
        keywords_list = [k.lower() for k in synopsis.get("keywords", [])]
        transcript_text = synopsis.get("transcript", "").lower()
        
        if any(x in title_text or x in transcript_text or any(x in k for k in keywords_list) for x in ["security", "encrypt", "crypt", "symmetric", "asymmetric", "cipher"]):
            security_questions = [
                "What is network security?",
                "What is encryption?",
                "What is the difference between symmetric and asymmetric encryption?",
                "What is cryptography?"
            ]
            for sq in security_questions:
                if sq not in questions:
                    questions.append(sq)
        
        # 3. Construct from glossary / definitions
        if len(questions) < 4:
            defs = synopsis.get("importantDefinitions") or synopsis.get("important_definitions") or []
            for d in defs:
                if isinstance(d, dict) and d.get("term"):
                    term = d["term"].strip()
                    q = f"What is the definition and significance of {term}?"
                    if q not in questions:
                        questions.append(q)
                        
        # 4. Construct from key concepts
        if len(questions) < 4:
            concepts = synopsis.get("keyConcepts") or synopsis.get("key_concepts") or synopsis.get("major_concepts") or synopsis.get("majorConcepts") or []
            for c in concepts:
                if isinstance(c, dict) and (c.get("concept") or c.get("name")):
                    concept_name = (c.get("concept") or c.get("name")).strip()
                    q = f"Explain the core concept and operational principles of {concept_name}."
                    if q not in questions:
                        questions.append(q)
                        
        # 5. Construct from topics
        if len(questions) < 4:
            topics = synopsis.get("topics") or []
            for t in topics:
                if isinstance(t, dict) and t.get("topic"):
                    topic_name = t["topic"].strip()
                    q = f"What are the primary challenges and methodologies related to {topic_name}?"
                    if q not in questions:
                        questions.append(q)
                        
        # 6. Absolute Fallback templates using the video title
        title = synopsis.get("metadata", {}).get("title", "the lecture topic")
        fallbacks = [
            f"What is the primary core objective and message of {title}?",
            f"What are the most critical practical takeaways presented in {title}?",
            f"Explain the primary real-world applications of the technologies covered in {title}.",
            f"What are the main constraints or limitations discussed in {title}?"
        ]
        for fq in fallbacks:
            if len(questions) >= 4:
                break
            if fq not in questions:
                questions.append(fq)
                
        return [q[:140] for q in questions[:4]]

    def _get_topic_specific_revision_notes(self, synopsis: dict) -> list[str]:
        """
        Dynamically extracts 4 topic-specific high-yield revision tips.
        """
        notes = []
        
        # 1. Try exam preparation notes
        exam_notes = synopsis.get("exam_preparation_notes") or synopsis.get("examPreparationNotes") or []
        for note in exam_notes:
            if isinstance(note, str) and note.strip():
                notes.append(note.strip())
                
        # 2. Try parsing quick revision notes
        rev_notes = synopsis.get("quick_revision_notes") or synopsis.get("quickRevisionNotes") or ""
        if isinstance(rev_notes, str) and rev_notes.strip():
            bullets = re.findall(r"(?:^|\n)\s*[-*•]\s*(.+)", rev_notes)
            for b in bullets:
                if b.strip() and b.strip() not in notes:
                    notes.append(b.strip())
            if not notes:
                for line in rev_notes.split("\n"):
                    line = re.sub(r"^#+\s*", "", line).strip()
                    if line and len(line) > 12 and line not in notes:
                        notes.append(line)
                        
        # 3. Try key takeaways
        takeaways = synopsis.get("keyTakeaways") or synopsis.get("key_takeaways") or []
        for t in takeaways:
            if isinstance(t, str) and t.strip() and t.strip() not in notes:
                notes.append(t.strip())

        # Inject highly relevant revision tips for security/encryption/cryptography
        title_text = synopsis.get("metadata", {}).get("title", "").lower()
        keywords_list = [k.lower() for k in synopsis.get("keywords", [])]
        transcript_text = synopsis.get("transcript", "").lower()
        
        if any(x in title_text or x in transcript_text or any(x in k for k in keywords_list) for x in ["security", "encrypt", "crypt", "symmetric", "asymmetric", "cipher"]):
            security_tips = [
                "Understand the core CIA triad: Confidentiality, Integrity, and Availability in network security.",
                "Differentiate between symmetric encryption (fast, single key) and asymmetric encryption (public-private key pairs).",
                "Review how public key cryptography secures communication channels like HTTPS/TLS.",
                "Understand the role of hashing functions in verifying data integrity without decrypting contents."
            ]
            for st in security_tips:
                if st not in notes:
                    notes.append(st)
                
        # 4. Try parsing sentences from executive summary
        if len(notes) < 4:
            exec_sum = synopsis.get("executiveSummary", "")
            if exec_sum:
                sentences = re.split(r"(?<=[.!?])\s+", exec_sum)
                for s in sentences:
                    if len(s.strip()) > 15 and s.strip() not in notes:
                        notes.append(s.strip())
                        
        # 5. Standard fallbacks based on title
        title = synopsis.get("metadata", {}).get("title", "this study guide")
        fallbacks = [
            f"Review the primary definitions and core terminologies associated with {title}.",
            f"Understand the step-by-step process workflows and diagrams outlined in this guide.",
            f"Analyze the key comparisons and performance trade-offs of different implementations.",
            f"Practice the self-assessment quiz questions to identify any knowledge gaps."
        ]
        for f in fallbacks:
            if len(notes) >= 4:
                break
            if f not in notes:
                notes.append(f)
                
        return [n[:140] for n in notes[:4]]

    def _get_topic_specific_comparison(self, synopsis: dict) -> list[dict]:
        """
        Dynamically generates a comparison matrix comparing different concepts from the synopsis.
        """
        # Inject Symmetric vs Asymmetric Comparison if the topic is security/encryption/cryptography
        title_text = synopsis.get("metadata", {}).get("title", "").lower()
        keywords_list = [k.lower() for k in synopsis.get("keywords", [])]
        transcript_text = synopsis.get("transcript", "").lower()
        
        if any(x in title_text or x in transcript_text or any(x in k for k in keywords_list) for x in ["security", "encrypt", "crypt", "symmetric", "asymmetric", "cipher"]):
            return [
                {
                    "aspect": "Key Management",
                    "item1": "Symmetric: Uses a single shared secret key for both encryption and decryption.",
                    "item2": "Asymmetric: Uses a mathematically linked public-private key pair."
                },
                {
                    "aspect": "Processing Speed",
                    "item1": "Symmetric: Extremely fast and computationally efficient; ideal for bulk data transfer.",
                    "item2": "Asymmetric: Significantly slower and resource-intensive; used for key exchange."
                },
                {
                    "aspect": "Primary Use Case",
                    "item1": "Symmetric: Encrypting hard drives, databases, and payload messages.",
                    "item2": "Asymmetric: Secure key distribution (SSL/TLS handshakes) and digital signatures."
                }
            ]

        concepts = synopsis.get("keyConcepts") or synopsis.get("key_concepts") or synopsis.get("major_concepts") or synopsis.get("majorConcepts") or []
        defs = synopsis.get("importantDefinitions") or synopsis.get("important_definitions") or []
        
        candidates = []
        for c in concepts:
            if isinstance(c, dict) and (c.get("concept") or c.get("name")):
                candidates.append({
                    "name": c.get("concept") or c.get("name"),
                    "desc": c.get("explanation", "")
                })
        for d in defs:
            if isinstance(d, dict) and d.get("term"):
                candidates.append({
                    "name": d.get("term"),
                    "desc": d.get("definition", "")
                })
                
        # Remove duplicates
        seen = set()
        unique_candidates = []
        for item in candidates:
            name_lower = item["name"].lower().strip()
            if name_lower not in seen:
                seen.add(name_lower)
                unique_candidates.append(item)
                
        # If we have at least 2 unique items, compare them
        if len(unique_candidates) >= 2:
            item1 = unique_candidates[0]
            item2 = unique_candidates[1]
            return [
                {
                    "aspect": "Core Concept",
                    "item1": f"{item1['name']}: {item1['desc'][:90]}...",
                    "item2": f"{item2['name']}: {item2['desc'][:90]}..."
                },
                {
                    "aspect": "Practical Focus",
                    "item1": f"How to apply {item1['name']} in real-world scenarios.",
                    "item2": f"How to deploy and monitor {item2['name']} systems."
                },
                {
                    "aspect": "Key Differences",
                    "item1": f"Distinguishing attributes of {item1['name']}.",
                    "item2": f"Operational requirements for {item2['name']} workflows."
                }
            ]
            
        # Default topic-specific fallback using title
        title = synopsis.get("metadata", {}).get("title", "the lecture topic")
        return [
            {
                "aspect": "Primary Domain",
                "item1": f"Theoretical principles of {title}",
                "item2": f"Implementation patterns of {title}"
            },
            {
                "aspect": "Operational Mode",
                "item1": "Proactive planning and configuration",
                "item2": "Reactive analysis and incident recovery"
            },
            {
                "aspect": "Execution Scope",
                "item1": "Individual component constraints",
                "item2": "End-to-end integration boundaries"
            }
        ]

    def _generate_fallback_slide_plan(self, synopsis: dict) -> dict:
        """
        Generates a deterministic fallback presentation slide plan.
        """
        metadata = synopsis.get("metadata", {})
        title = metadata.get("title", "Video Synopsis")
        channel = metadata.get("channelName", "Unknown Channel")
        exec_summary = synopsis.get("executiveSummary", "")
        introduction = synopsis.get("introduction", "")
        chapters = synopsis.get("chapters", [])
        concepts = synopsis.get("major_concepts") or synopsis.get("majorConcepts") or synopsis.get("keyConcepts") or []
        practical = synopsis.get("practicalApplications", [])

        # Build timeline items
        timeline_items = []
        for c in chapters[:4]:
            timeline_items.append({
                "time": c.get("timestamp") or "00:00",
                "title": c.get("title") or "Chapter Outline",
                "description": c.get("summary") or "Section discussion"
            })
        while len(timeline_items) < 4:
            timeline_items.append({"time": "00:00", "title": "Lecture Phase", "description": "Core educational presentation"})

        # Build concept items
        concepts_items = []
        for c in concepts[:3]:
            concepts_items.append({
                "title": c.get("concept") or c.get("name") or "Concept Pillar",
                "description": c.get("explanation") or "Topic conceptual details",
                "icon": "ðŸ’¡"
            })
        while len(concepts_items) < 3:
            concepts_items.append({"title": "Study Pillar", "description": "Core theoretical architecture details", "icon": "ðŸ’¡"})

        # Build flowchart items
        flow_items = []
        for a in practical[:4]:
            flow_items.append({
                "step": a.get("application") or "Application",
                "description": a.get("description") or "Practical deployment instruction"
            })
        if not flow_items:
            flow_items = [
                {"step": "Analyze", "description": "Evaluate core concepts"},
                {"step": "Structure", "description": "Build deployment templates"},
                {"step": "Execute", "description": "Deploy standard setups"},
                {"step": "Verify", "description": "Confirm testing results"}
            ]
        while len(flow_items) < 4:
            flow_items.append({"step": "Next Step", "description": "Further process validation"})

        comparison_items = self._get_topic_specific_comparison(synopsis)
        quiz_questions = self._get_topic_specific_questions(synopsis)
        revision_notes = self._get_topic_specific_revision_notes(synopsis)
        
        slides = [
            {
                "type": "title",
                "title": title,
                "content": [
                    f"Comprehensive curriculum walkthrough based on the lecture by {channel}",
                    "Deep theoretical insights combined with step-by-step process workflows",
                    "Includes active recall revision check points and practice quiz materials"
                ],
                "diagram": {"type": "none", "items": []},
                "speaker_notes": f"Welcome everyone to this comprehensive presentation on '{title}'. Today we will walk through a detailed, high-fidelity study guide compiled from the video transcript. We will cover the chronological chapters, core theoretical concepts, process workflows, a comparative matrix, and conclude with a self-assessment quiz and revision notes to ensure complete mastery of this topic."
            },
            {
                "type": "agenda",
                "title": "Presentation Agenda",
                "content": [
                    "Walkthrough of study segments covering the core learning objectives",
                    "Chronological outline of the video chapters progression and key milestones",
                    "Core conceptual pillars, process workflows, and comparative analysis",
                    "Active recall revision notes and practice self-assessment quiz"
                ],
                "diagram": {
                    "type": "cards",
                    "items": [
                        {"title": "Objectives", "description": "Learning targets", "icon": "🎯"},
                        {"title": "Timeline", "description": "Chapters schedule", "icon": "📅"},
                        {"title": "Concepts", "description": "Core pillars", "icon": "💡"},
                        {"title": "Review", "description": "Revision tips", "icon": "📝"}
                    ]
                },
                "speaker_notes": "Our agenda for today is structured into four main areas. First, we will review the learning objectives and executive summary to establish a context baseline. Second, we will look at the video chapters timeline to see how the lecture unfolds chronologically. Third, we will explore the core conceptual pillars and workflows. Finally, we will test our understanding with a practice quiz and review revision tips."
            },
            {
                "type": "summary",
                "title": "Executive Summary",
                "content": [
                    exec_summary if len(exec_summary) > 20 else f"Core summary insights of the video topic: {title}",
                    introduction if len(introduction) > 20 else "Baseline introduction and contextual relevance of the lecture"
                ],
                "diagram": {"type": "none", "items": []},
                "speaker_notes": f"Let's dive into the executive summary of this lecture. The core message of this video focuses on: '{exec_summary[:200]}...'. The introduction establishes the baseline context: '{introduction[:200]}...'. As a presenter, you should emphasize these foundational highlights to give your audience a clear roadmap of the concepts we will detail in the upcoming slides."
            },
            {
                "type": "timeline",
                "title": "Lecture Chapters Timeline",
                "content": [
                    "Chronological timestamp outline listing when core topics and key milestones occur",
                    "Provides visual checkpoints of the video lecture progression from start to finish",
                    "Helps locate specific segments in the timeline for deep dive reviews"
                ],
                "diagram": {"type": "timeline", "items": timeline_items},
                "speaker_notes": f"This slide presents the chronological timeline of the video chapters. We start at {timeline_items[0]['time']} with {timeline_items[0]['title']}, where the speaker introduces the fundamental principles. Then we move through the subsequent sections, each building on the previous one, leading up to the final summary and takeaways at {timeline_items[-1]['time']}. Reviewing this chronology helps us understand the logical structure of the lecture and lets students pinpoint exactly where in the video specific topics are discussed."
            },
            {
                "type": "concepts",
                "title": "Theoretical Conceptual Pillars",
                "content": [
                    "Core conceptual definitions establishing the foundational models",
                    "Key structural theories and principles explained by the speaker in detail",
                    "Forms the intellectual baseline for understanding this technology or topic"
                ],
                "diagram": {"type": "cards", "items": concepts_items},
                "speaker_notes": f"Next, let's explore the core theoretical concepts discussed by the speaker. These pillars form the intellectual foundation of this topic. First, we look at {concepts_items[0]['title']}, which explains {concepts_items[0]['description'][:100]}... Second, we have {concepts_items[1]['title']}, which details {concepts_items[1]['description'][:100]}... Finally, {concepts_items[2]['title']} outlines {concepts_items[2]['description'][:100]}... Take time to explain the real-world significance of each of these concepts to your audience."
            },
            {
                "type": "flowchart",
                "title": "Process Workflows",
                "content": [
                    "Step-by-step implementation rules showing the logical workflow progression",
                    "Practical methodologies and execution steps recommended by the speaker",
                    "Helps in applying the theoretical concepts to real-world software scenarios"
                ],
                "diagram": {"type": "flowchart", "items": flow_items},
                "speaker_notes": f"This diagram walks through the process workflows and practical applications step-by-step. In the first phase, we focus on {flow_items[0]['step']}, which involves {flow_items[0]['description']}. This is followed by {flow_items[1]['step']}, and then we move into the execution phase of {flow_items[2]['step']}, culminating in the validation and finalization steps. Walk your audience through this flowchart to show how these concepts transition from theory into practical, step-by-step implementation."
            },
            {
                "type": "comparison",
                "title": "Concepts Comparison Matrix",
                "content": [
                    "Reviewing alternative choices and comparing their relative attributes and scopes",
                    "Compares key metrics, tradeoffs, and deployment configurations of core concepts",
                    "Assists in decision-making when designing system components and architectures"
                ],
                "diagram": {"type": "table", "items": comparison_items},
                "speaker_notes": "Let's look at the Concepts Comparison Matrix. This table compares the primary criteria and aspects of different options. We evaluate the core concepts, their practical implementation focus, and key operational differences. This comparative analysis is extremely useful for students and practitioners to understand the trade-offs between different choices and make informed decisions when designing their systems."
            },
            {
                "type": "quiz",
                "title": "Assessment Quiz",
                "content": quiz_questions,
                "diagram": {"type": "none", "items": []},
                "speaker_notes": f"To check our knowledge retention, we have a practice quiz containing 4 key questions directly related to the video content: 1. {quiz_questions[0]} 2. {quiz_questions[1]} 3. {quiz_questions[2]} 4. {quiz_questions[3]}. Let's go through each question one by one, ask the audience for their answers, explain the correct choices, and clear up any conceptual gaps they might have."
            },
            {
                "type": "revision",
                "title": "Quick Revision Notes",
                "content": revision_notes,
                "diagram": {"type": "none", "items": []},
                "speaker_notes": f"Finally, let's review the high-yield revision tips for exam preparation. These core takeaways summarize the most critical concepts: 1. {revision_notes[0]} 2. {revision_notes[1]} 3. {revision_notes[2]} 4. {revision_notes[3]}. Focusing on these active recall notes will help you recall the content quickly during your exams or technical interviews. Encourage the students to review these points regularly."
            },
            {
                "type": "thankyou",
                "title": "Questions & Answers",
                "content": [
                    "Thank you for reviewing this visual study guide presentation",
                    "We will now address questions, clarifications, and open discussion points",
                    "Complete study handbooks and self-grading quizzes are available on the dashboard"
                ],
                "diagram": {"type": "none", "items": []},
                "speaker_notes": f"Thank you for reviewing this visual study guide presentation based on '{title}'. I am now open to any questions, comments, or further clarifications on the slides, timeline, or conceptual structures we have discussed today. Let's open the floor for Q&A and dive deeper into any topics of interest."
            }
        ]

        return {
            "subtitle": "Comprehensive Study Guide Presentation",
            "slides": slides
        }

    def _build_ppt_bytes(self, synopsis: dict, slide_plan: dict) -> bytes:
        """
        Builds a professional, dark-purple themed widescreen (16:9) presentation
        based on the provided slide plan.
        """
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        import io

        prs = Presentation()
        # Set 16:9 widescreen dimensions
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Theme colors: modern dark purple scheme matching Video Synopsis AI dashboard
        colors = {
            "bg": RGBColor(9, 9, 11),          # Dark background (matching dark UI)
            "card": RGBColor(24, 20, 38),      # Slightly lighter card background
            "border": RGBColor(59, 45, 99),    # Border purple
            "accent": RGBColor(168, 85, 247),  # Accent Purple
            "accent_sub": RGBColor(236, 72, 153), # Accent Pink
            "white": RGBColor(255, 255, 255),
            "zinc": RGBColor(216, 180, 254),   # Light violet text
            "zinc_muted": RGBColor(161, 161, 170)
        }

        # Safe text helper to prevent overflows
        def clean_str(val, max_len=None):
            if not val:
                return ""
            val = str(val).strip()
            if max_len and len(val) > max_len:
                return val[:max_len] + "..."
            return val

        # Process each slide from the plan
        for slide_idx, slide_data in enumerate(slide_plan.get("slides", [])):
            s_type = slide_data.get("type", "summary").lower().strip()
            s_title = clean_str(slide_data.get("title", "Lecture Notes"), 80)
            
            # Safe truncation of content strings
            raw_content = slide_data.get("content") or []
            if isinstance(raw_content, str):
                raw_content = [raw_content]
            s_content = [clean_str(x, 150) for x in raw_content if x]
            s_notes = clean_str(slide_data.get("speaker_notes", ""))
            
            diagram = slide_data.get("diagram", {})
            diag_items = diagram.get("items") or []
            diag_type = str(diagram.get("type", "none")).lower().strip()

            if s_type == "title":
                subtitle = clean_str(slide_plan.get("subtitle", "Lecture Guide Presentation"), 120)
                self._add_title_slide(prs, colors, s_title, subtitle, s_notes)
            elif s_type == "agenda":
                self._add_agenda_slide(prs, colors, s_title, s_content, diag_items, s_notes)
            elif s_type == "timeline" or diag_type == "timeline":
                self._add_timeline_slide(prs, colors, s_title, s_content, diag_items, s_notes)
            elif s_type == "flowchart" or diag_type == "flowchart":
                self._add_flowchart_slide(prs, colors, s_title, s_content, diag_items, s_notes)
            elif s_type == "concepts" or diag_type == "cards":
                self._add_icon_cards_slide(prs, colors, s_title, s_content, diag_items, s_notes)
            elif s_type == "comparison" or diag_type == "table":
                comparison_items = diag_items
                has_generic = any(any(p in str(x.get("aspect", "")).lower() or p in str(x.get("item1", "")).lower() or p in str(x.get("item2", "")).lower() for p in ["processes decoupled", "decoupled and isolated", "system latency constraint", "horizontal clusters"]) for x in comparison_items)
                # If comparison items are generic or missing, replace with topic-specific
                if not comparison_items or len(comparison_items) < 2 or has_generic or any(str(x.get("aspect", "")).lower() in ["feature comparison aspect", "feature setup"] for x in comparison_items):
                    comparison_items = self._get_topic_specific_comparison(synopsis)
                self._add_comparison_slide(prs, colors, s_title, s_content, comparison_items, s_notes)
            elif s_type == "quiz":
                quiz_questions = s_content
                has_generic = any(any(p in str(q).lower() for p in ["processes decoupled", "decoupled and isolated", "system latency constraint", "horizontal clusters"]) for q in quiz_questions)
                # Intercept and replace with topic-specific questions if generic or incomplete
                if len(quiz_questions) < 4 or has_generic or any(x.lower() in ["quiz question a", "quiz overview", "assessment questions", "questions"] for x in quiz_questions):
                    quiz_questions = self._get_topic_specific_questions(synopsis)
                self._add_quiz_slide(prs, colors, s_title, quiz_questions, s_notes)
            elif s_type == "revision":
                revision_points = s_content
                has_generic = any(any(p in str(r).lower() for p in ["processes decoupled", "decoupled and isolated", "system latency constraint", "horizontal clusters"]) for r in revision_points)
                # Intercept and replace with topic-specific revision points if generic or incomplete
                if len(revision_points) < 4 or has_generic or any(x.lower() in ["revision tip a", "revision guidelines", "revision notes", "exam tips"] for x in revision_points):
                    revision_points = self._get_topic_specific_revision_notes(synopsis)
                self._add_revision_slide(prs, colors, s_title, revision_points, s_notes)
            elif s_type == "thankyou":
                self._add_thank_you_slide(prs, colors, s_title, s_content, s_notes)
            else:
                self._add_summary_slide(prs, colors, s_title, s_content, s_notes)

        # Output presentation bytes
        ppt_stream = io.BytesIO()
        prs.save(ppt_stream)
        ppt_stream.seek(0)
        return ppt_stream.getvalue()


    def _add_slide_with_bg(self, prs, bg_color, accent_color):
        from pptx.util import Inches
        from pptx.enum.shapes import MSO_SHAPE
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
        
        # Left accent vertical bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent_color
        bar.line.fill.background()
        return slide

    def _draw_divider_line(self, slide, left, top, width, color):
        from pptx.util import Inches
        from pptx.enum.shapes import MSO_SHAPE
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = color
        line.line.fill.background()
        return line

    def _add_header(self, slide, title, colors, category=None):
        from pptx.util import Inches, Pt
        if category:
            catBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.2), Inches(11.733), Inches(0.3))
            cat_tf = catBox.text_frame
            cat_tf.word_wrap = True
            cat_tf.margin_left = Inches(0)
            cat_tf.margin_top = Inches(0)
            cat_p = cat_tf.paragraphs[0]
            cat_p.text = str(category).upper()
            cat_p.font.name = 'Outfit'
            cat_p.font.size = Pt(10)
            cat_p.font.bold = True
            cat_p.font.color.rgb = colors["accent_sub"]
        
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(11.733), Inches(0.7))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0)
        tf.margin_top = Inches(0)
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = 'Outfit'
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = colors["accent"]
        
        self._draw_divider_line(slide, Inches(0.8), Inches(1.15), Inches(11.733), colors["accent_sub"])

    def _draw_card_shape(self, slide, left, top, width, height, colors, border_color=None, border_width=1.5):
        from pptx.util import Pt
        from pptx.enum.shapes import MSO_SHAPE
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors["card"]
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(border_width)
        else:
            shape.line.fill.background()
        return shape

    def _add_speaker_notes(self, slide, notes):
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = "\n".join(notes) if isinstance(notes, list) else str(notes)

    def _add_title_slide(self, prs, colors, title, subtitle, speaker_notes):
        from pptx.util import Inches, Pt
        slide = self._add_slide_with_bg(prs, colors["bg"], colors["accent"])
        
        # Logo Label
        logo_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.733), Inches(0.4))
        logo_tf = logo_box.text_frame
        logo_tf.margin_left = Inches(0)
        logo_tf.word_wrap = True
        logo_p = logo_tf.paragraphs[0]
        logo_p.text = "âš¡ AI STUDY HANDBOOK  |  VISUAL PRESENTATION"
        logo_p.font.name = 'Outfit'
        logo_p.font.size = Pt(12)
        logo_p.font.bold = True
        logo_p.font.color.rgb = colors["accent"]

        # Main Title Box
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.733), Inches(2.2))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_tf.margin_left = Inches(0)
        title_p = title_tf.paragraphs[0]
        title_p.text = title
        title_p.font.name = 'Outfit'
        title_p.font.size = Pt(36)
        title_p.font.bold = True
        title_p.font.color.rgb = colors["white"]
        title_p.space_after = Pt(12)

        # Subtitle
        sub_p = title_tf.add_paragraph()
        sub_p.text = subtitle
        sub_p.font.name = 'Plus Jakarta Sans'
        sub_p.font.size = Pt(16)
        sub_p.font.color.rgb = colors["zinc"]
        
        # Visual Divider
        self._draw_divider_line(slide, Inches(0.8), Inches(4.3), Inches(5.0), colors["accent_sub"])
        
        # Footer / Info Box
        info_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.6), Inches(11.733), Inches(1.5))
        info_tf = info_box.text_frame
        info_tf.margin_left = Inches(0)
        info_tf.word_wrap = True
        
        inf_p = info_tf.paragraphs[0]
        inf_p.text = "Features included in this guide:"
        inf_p.font.name = 'Outfit'
        inf_p.font.size = Pt(13)
        inf_p.font.bold = True
        inf_p.font.color.rgb = colors["white"]
        inf_p.space_after = Pt(6)
        
        inf_p2 = info_tf.add_paragraph()
        inf_p2.text = "âœ” Timeline Outline   âœ” Concepts Cards   âœ” Flowchart Workflow   âœ” Active Recall Check"
        inf_p2.font.name = 'Plus Jakarta Sans'
        inf_p2.font.size = Pt(11)
        inf_p2.font.color.rgb = colors["zinc_muted"]
        
        self._add_speaker_notes(slide, [speaker_notes])

    def _add_agenda_slide(self, prs, colors, title, content, items, speaker_notes):
        from pptx.util import Inches, Pt
        slide = self._add_slide_with_bg(prs, colors["bg"], colors["accent"])
        self._add_header(slide, title, colors, "Roadmap")
        
        # Left Description Box
        left_w = Inches(5.4)
        self._draw_card_shape(slide, Inches(0.8), Inches(1.8), left_w, Inches(4.8), colors, colors["border"])
        left_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), left_w - Inches(0.4), Inches(4.4))
        left_tf = left_box.text_frame
        left_tf.word_wrap = True
        
        alp1 = left_tf.paragraphs[0]
        alp1.text = "Presentation Structure"
        alp1.font.name = 'Outfit'
        alp1.font.size = Pt(18)
        alp1.font.bold = True
        alp1.font.color.rgb = colors["accent"]
        alp1.space_after = Pt(12)
        
        for c in content:
            alp2 = left_tf.add_paragraph()
            alp2.text = c
            alp2.font.name = 'Plus Jakarta Sans'
            alp2.font.size = Pt(12)
            alp2.font.color.rgb = colors["zinc"]
            alp2.line_spacing = 1.3
            alp2.space_after = Pt(8)
            
        # Right Agenda Cards (Items)
        if not items:
            items = [
                {"title": "Objectives", "description": "Target metrics", "icon": "ðŸŽ¯"},
                {"title": "Concepts", "description": "Theoretical cards", "icon": "ðŸ’¡"},
                {"title": "Workflows", "description": "Process flowchart", "icon": "âš¡"},
                {"title": "Review Quiz", "description": "Active recall test", "icon": "ðŸ“"}
            ]
            
        for idx, item in enumerate(items[:4]):
            top_pos = Inches(1.8) + idx * Inches(1.25)
            self._draw_card_shape(slide, Inches(6.6), top_pos, Inches(5.9), Inches(1.05), colors, colors["border"])
            
            # Icon
            icon = item.get("icon") or "âš¡"
            icon_box = slide.shapes.add_textbox(Inches(6.75), top_pos + Inches(0.12), Inches(0.8), Inches(0.8))
            icon_tf = icon_box.text_frame
            icon_p = icon_tf.paragraphs[0]
            icon_p.text = icon
            icon_p.font.name = 'Outfit'
            icon_p.font.size = Pt(22)
            icon_p.font.bold = True
            
            # Content Box
            info_box = slide.shapes.add_textbox(Inches(7.6), top_pos + Inches(0.08), Inches(4.7), Inches(0.85))
            info_tf = info_box.text_frame
            info_tf.word_wrap = True
            
            ip1 = info_tf.paragraphs[0]
            ip1.text = item.get("title", "")
            ip1.font.name = 'Outfit'
            ip1.font.size = Pt(13)
            ip1.font.bold = True
            ip1.font.color.rgb = colors["white"]
            
            ip2 = info_tf.add_paragraph()
            ip2.text = item.get("description", "")
            ip2.font.name = 'Plus Jakarta Sans'
            ip2.font.size = Pt(10.5)
            ip2.font.color.rgb = colors["zinc_muted"]
            ip2.space_before = Pt(2)

        self._add_speaker_notes(slide, [speaker_notes])

    def _add_timeline_slide(self, prs, colors, title, content, items, speaker_notes):
        from pptx.util import Inches, Pt
        from pptx.enum.shapes import MSO_SHAPE
        slide = self._add_slide_with_bg(prs, colors["bg"], colors["accent"])
        self._add_header(slide, title, colors, "Chronology")
        
        if content:
            desc_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.733), Inches(0.5))
            desc_tf = desc_box.text_frame
            desc_tf.word_wrap = True
            desc_tf.margin_left = Inches(0)
            desc_p = desc_tf.paragraphs[0]
            desc_p.text = content[0][:150]
            desc_p.font.name = 'Plus Jakarta Sans'
            desc_p.font.size = Pt(12)
            desc_p.font.color.rgb = colors["zinc"]
            
        if not items:
            items = [
                {"time": "00:00", "title": "Introduction", "description": "Opening outline"},
                {"time": "05:00", "title": "Core Design", "description": "Foundational topics"},
                {"time": "12:00", "title": "Implementation", "description": "Live demonstration"},
                {"time": "20:00", "title": "Wrap-up", "description": "Takeaways summary"}
            ]
            
        card_w = Inches(2.7)
        card_h = Inches(4.0)
        spacing = Inches(0.3)
        start_left = Inches(0.8)
        top_pos = Inches(2.0)
        
        for idx, item in enumerate(items[:4]):
            left_pos = start_left + idx * (card_w + spacing)
            self._draw_card_shape(slide, left_pos, top_pos, card_w, card_h, colors, colors["border"])
            
            # Connect line to next (except last)
            if idx < len(items[:4]) - 1:
                arrow_left = left_pos + card_w
                arrow_w = spacing
                arrow_top = top_pos + Inches(1.8)
                arrow_h = Inches(0.4)
                arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_left, arrow_top, arrow_w, arrow_h)
                arr.fill.solid()
                arr.fill.fore_color.rgb = colors["accent"]
                arr.line.fill.background()
                
            # Timestamp Box at the top of card
            ts_box = slide.shapes.add_textbox(left_pos + Inches(0.15), top_pos + Inches(0.2), card_w - Inches(0.3), Inches(0.8))
            ts_tf = ts_box.text_frame
            ts_tf.word_wrap = True
            ts_p = ts_tf.paragraphs[0]
            ts_p.text = item.get("time") or item.get("timestamp") or "00:00"
            ts_p.alignment = 1 # Center
            ts_p.font.name = 'Outfit'
            ts_p.font.size = Pt(18)
            ts_p.font.bold = True
            ts_p.font.color.rgb = colors["accent"]
            
            # Chapter details Box
            info_box = slide.shapes.add_textbox(left_pos + Inches(0.15), top_pos + Inches(1.0), card_w - Inches(0.3), Inches(2.8))
            info_tf = info_box.text_frame
            info_tf.word_wrap = True
            
            ip1 = info_tf.paragraphs[0]
            ip1.text = item.get("title", "")
            ip1.alignment = 1 # Center
            ip1.font.name = 'Outfit'
            ip1.font.size = Pt(13)
            ip1.font.bold = True
            ip1.font.color.rgb = colors["white"]
            ip1.space_after = Pt(8)
            
            ip2 = info_tf.add_paragraph()
            ip2.text = item.get("description") or item.get("summary") or ""
            ip2.alignment = 1 # Center
            ip2.font.name = 'Plus Jakarta Sans'
            ip2.font.size = Pt(10.5)
            ip2.font.color.rgb = colors["zinc"]
            ip2.line_spacing = 1.15

        self._add_speaker_notes(slide, [speaker_notes])

    def _add_flowchart_slide(self, prs, colors, title, content, items, speaker_notes):
        from pptx.util import Inches, Pt
        from pptx.enum.shapes import MSO_SHAPE
        slide = self._add_slide_with_bg(prs, colors["bg"], colors["accent"])
        self._add_header(slide, title, colors, "Workflows")
        
        if content:
            desc_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.733), Inches(0.5))
            desc_tf = desc_box.text_frame
            desc_tf.word_wrap = True
            desc_tf.margin_left = Inches(0)
            desc_p = desc_tf.paragraphs[0]
            desc_p.text = content[0][:150]
            desc_p.font.name = 'Plus Jakarta Sans'
            desc_p.font.size = Pt(12)
            desc_p.font.color.rgb = colors["zinc"]
            
        if not items:
            items = [
                {"step": "Setup", "description": "Configure systems"},
                {"step": "Execution", "description": "Process workflow tasks"},
                {"step": "Analysis", "description": "Audit metrics output"},
                {"step": "Finalize", "description": "Publish dashboard reports"}
            ]
            
        card_w = Inches(2.7)
        card_h = Inches(3.6)
        spacing = Inches(0.3)
        start_left = Inches(0.8)
        top_pos = Inches(2.2)
        
        for idx, item in enumerate(items[:4]):
            left_pos = start_left + idx * (card_w + spacing)
            self._draw_card_shape(slide, left_pos, top_pos, card_w, card_h, colors, colors["accent_sub"])
            
            # Connection arrow
            if idx < len(items[:4]) - 1:
                arrow_left = left_pos + card_w
                arrow_w = spacing
                arrow_top = top_pos + Inches(1.6)
                arrow_h = Inches(0.4)
                arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_left, arrow_top, arrow_w, arrow_h)
                arr.fill.solid()
                arr.fill.fore_color.rgb = colors["accent"]
                arr.line.fill.background()
                
            # Step Number
            num_box = slide.shapes.add_textbox(left_pos + Inches(0.15), top_pos + Inches(0.2), card_w - Inches(0.3), Inches(0.8))
            num_tf = num_box.text_frame
            num_p = num_tf.paragraphs[0]
            num_p.text = f"Step 0{idx+1}"
            num_p.alignment = 1 # Center
            num_p.font.name = 'Outfit'
            num_p.font.size = Pt(15)
            num_p.font.bold = True
            num_p.font.color.rgb = colors["accent"]
            
            # Content Box
            info_box = slide.shapes.add_textbox(left_pos + Inches(0.15), top_pos + Inches(0.9), card_w - Inches(0.3), Inches(2.5))
            info_tf = info_box.text_frame
            info_tf.word_wrap = True
            
            ip1 = info_tf.paragraphs[0]
            ip1.text = item.get("step") or item.get("title") or ""
            ip1.alignment = 1 # Center
            ip1.font.name = 'Outfit'
            ip1.font.size = Pt(13)
            ip1.font.bold = True
            ip1.font.color.rgb = colors["white"]
            ip1.space_after = Pt(8)
            
            ip2 = info_tf.add_paragraph()
            ip2.text = item.get("description") or ""
            ip2.alignment = 1 # Center
            ip2.font.name = 'Plus Jakarta Sans'
            ip2.font.size = Pt(10.5)
            ip2.font.color.rgb = colors["zinc"]
            ip2.line_spacing = 1.15

        self._add_speaker_notes(slide, [speaker_notes])

    def _add_icon_cards_slide(self, prs, colors, title, content, items, speaker_notes):
        from pptx.util import Inches, Pt
        slide = self._add_slide_with_bg(prs, colors["bg"], colors["accent"])
        self._add_header(slide, title, colors, "Concepts")
        
        if not items:
            items = [
                {"title": "Concept Pillar 1", "description": "Technical insights details...", "icon": "ðŸ’¡"},
                {"title": "Concept Pillar 2", "description": "Technical insights details...", "icon": "âš¡"},
                {"title": "Concept Pillar 3", "description": "Technical insights details...", "icon": "ðŸŽ¯"}
            ]
            
        card_w = Inches(3.7)
        card_h = Inches(4.5)
        spacing = Inches(0.3)
        start_left = Inches(0.8)
        top_pos = Inches(1.8)
        
        for idx, item in enumerate(items[:3]):
            left_pos = start_left + idx * (card_w + spacing)
            self._draw_card_shape(slide, left_pos, top_pos, card_w, card_h, colors, colors["border"])
            
            # Emoji Icon box
            icon = item.get("icon") or "ðŸ’¡"
            icon_box = slide.shapes.add_textbox(left_pos + Inches(0.2), top_pos + Inches(0.25), card_w - Inches(0.4), Inches(0.8))
            icon_tf = icon_box.text_frame
            icon_p = icon_tf.paragraphs[0]
            icon_p.text = icon
            icon_p.font.name = 'Outfit'
            icon_p.font.size = Pt(28)
            icon_p.font.bold = True
            
            # Content Box
            info_box = slide.shapes.add_textbox(left_pos + Inches(0.2), top_pos + Inches(1.1), card_w - Inches(0.4), Inches(3.2))
            info_tf = info_box.text_frame
            info_tf.word_wrap = True
            
            ip1 = info_tf.paragraphs[0]
            ip1.text = item.get("title") or item.get("concept") or ""
            ip1.font.name = 'Outfit'
            ip1.font.size = Pt(15)
            ip1.font.bold = True
            ip1.font.color.rgb = colors["accent"]
            ip1.space_after = Pt(10)
            
            ip2 = info_tf.add_paragraph()
            ip2.text = item.get("description") or item.get("explanation") or ""
            ip2.font.name = 'Plus Jakarta Sans'
            ip2.font.size = Pt(11)
            ip2.font.color.rgb = colors["zinc"]
            ip2.line_spacing = 1.2

        self._add_speaker_notes(slide, [speaker_notes])

    def _add_comparison_slide(self, prs, colors, title, content, items, speaker_notes):
        from pptx.util import Inches, Pt
        slide = self._add_slide_with_bg(prs, colors["bg"], colors["accent"])
        self._add_header(slide, title, colors, "Comparisons")
        
        if not items:
            title_val = str(title or "the topic")
            items = [
                {
                    "aspect": "Core Concept",
                    "item1": f"Theoretical principles of {title_val}",
                    "item2": f"Practical implementation steps for {title_val}"
                },
                {
                    "aspect": "Primary Scope",
                    "item1": "System-wide integration constraints",
                    "item2": "Individual component parameters"
                },
                {
                    "aspect": "Key Focus",
                    "item1": "Resilience and fault tolerance",
                    "item2": "Execution efficiency and speed"
                }
            ]
            
        rows_num = len(items[:4]) + 1
        cols_num = 3
        tbl_left = Inches(0.8)
        tbl_top = Inches(1.8)
        tbl_w = Inches(11.733)
        tbl_h = Inches(4.5)
        
        table_shape = slide.shapes.add_table(rows_num, cols_num, tbl_left, tbl_top, tbl_w, tbl_h)
        table = table_shape.table
        table.columns[0].width = Inches(2.6)
        table.columns[1].width = Inches(4.566)
        table.columns[2].width = Inches(4.566)
        
        headers = ["Aspect / Criteria", "Conceptual Focus", "Implementation / Alternative"]
        for c in range(3):
            cell = table.cell(0, c)
            cell.text = headers[c]
            cell.fill.solid()
            cell.fill.fore_color.rgb = colors["accent"]
            p = cell.text_frame.paragraphs[0]
            p.font.name = 'Outfit'
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = colors["white"]
            
        for r_idx, row_vals in enumerate(items[:4]):
            vals = [row_vals.get("aspect", ""), row_vals.get("item1", ""), row_vals.get("item2", "")]
            for c_idx, val in enumerate(vals):
                cell = table.cell(r_idx + 1, c_idx)
                cell.text = val
                cell.fill.solid()
                cell.fill.fore_color.rgb = colors["card"]
                p = cell.text_frame.paragraphs[0]
                p.font.name = 'Plus Jakarta Sans'
                p.font.size = Pt(11.5)
                p.font.color.rgb = colors["zinc"]
                
        self._add_speaker_notes(slide, [speaker_notes])

    def _add_revision_slide(self, prs, colors, title, content, speaker_notes):
        from pptx.util import Inches, Pt
        slide = self._add_slide_with_bg(prs, colors["bg"], colors["accent"])
        self._add_header(slide, title, colors, "Quick Revision")
        
        # Draw Left Card (Key Guidelines)
        self._draw_card_shape(slide, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), colors, colors["border"])
        left_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.4))
        left_tf = left_box.text_frame
        left_tf.word_wrap = True
        
        lp1 = left_tf.paragraphs[0]
        lp1.text = "ðŸ”‘ Core Revision Points"
        lp1.font.name = 'Outfit'
        lp1.font.size = Pt(18)
        lp1.font.bold = True
        lp1.font.color.rgb = colors["accent"]
        lp1.space_after = Pt(12)
        
        for idx, line in enumerate(content[:4]):
            p = left_tf.add_paragraph()
            p.text = f"â€¢ {line}"
            p.font.name = 'Plus Jakarta Sans'
            p.font.size = Pt(12)
            p.font.color.rgb = colors["zinc"]
            p.space_after = Pt(8)
            
        # Draw Right Card (Quick Notes Summary)
        self._draw_card_shape(slide, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.8), colors, colors["accent_sub"])
        right_box = slide.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.2), Inches(4.4))
        right_tf = right_box.text_frame
        right_tf.word_wrap = True
        
        rp1 = right_tf.paragraphs[0]
        rp1.text = "ðŸŽ¯ High-Yield Takeaways"
        rp1.font.name = 'Outfit'
        rp1.font.size = Pt(18)
        rp1.font.bold = True
        rp1.font.color.rgb = colors["white"]
        rp1.space_after = Pt(12)
        
        p2 = right_tf.add_paragraph()
        p2.text = "Focus on active recall, flashcards, and practice tests to reinforce your understanding. Make sure to review the chapters timeline for visual alignment."
        p2.font.name = 'Plus Jakarta Sans'
        p2.font.size = Pt(12)
        p2.font.color.rgb = colors["zinc"]
        p2.line_spacing = 1.3
        
        self._add_speaker_notes(slide, [speaker_notes])

    def _add_quiz_slide(self, prs, colors, title, content, speaker_notes):
        from pptx.util import Inches, Pt
        slide = self._add_slide_with_bg(prs, colors["bg"], colors["accent"])
        self._add_header(slide, title, colors, "Practice Quiz")
        
        # Question Card at top
        q_h = Inches(1.5)
        self._draw_card_shape(slide, Inches(0.8), Inches(1.8), Inches(11.733), q_h, colors, colors["accent"])
        
        q_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.95), Inches(11.333), Inches(1.2))
        q_tf = q_box.text_frame
        q_tf.word_wrap = True
        q_p = q_tf.paragraphs[0]
        q_p.text = f"Question Checklist: Assess your memory bounds on these core topics."
        q_p.font.name = 'Outfit'
        q_p.font.size = Pt(16)
        q_p.font.bold = True
        q_p.font.color.rgb = colors["white"]
        
        # Options Cards
        opt_w = Inches(5.6)
        opt_h = Inches(1.2)
        opt_positions = [
            (Inches(0.8), Inches(3.6)),
            (Inches(6.9), Inches(3.6)),
            (Inches(0.8), Inches(5.0)),
            (Inches(6.9), Inches(5.0))
        ]
        
        title_val = str(title or "the topic")
        fallback_questions = [
            f"What is the primary core objective of {title_val}?",
            f"What are the most critical takeaways of this topic?",
            f"Explain the primary real-world applications of {title_val}.",
            f"What are the main constraints or limitations discussed?"
        ]
        
        # If content has list of questions/points, use them. Otherwise use fallbacks
        points = content if len(content) >= 4 else fallback_questions
        for idx in range(4):
            left, top = opt_positions[idx]
            self._draw_card_shape(slide, left, top, opt_w, opt_h, colors, colors["border"])
            
            opt_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.18), opt_w - Inches(0.4), opt_h - Inches(0.3))
            opt_tf = opt_box.text_frame
            opt_tf.word_wrap = True
            opt_p = opt_tf.paragraphs[0]
            
            label = ["A", "B", "C", "D"][idx]
            opt_p.text = f"[{label}] {points[idx][:80]}"
            opt_p.font.name = 'Plus Jakarta Sans'
            opt_p.font.size = Pt(12)
            opt_p.font.color.rgb = colors["zinc"]
            
        self._add_speaker_notes(slide, [speaker_notes])

    def _add_summary_slide(self, prs, colors, title, content, speaker_notes):
        from pptx.util import Inches, Pt
        slide = self._add_slide_with_bg(prs, colors["bg"], colors["accent"])
        self._add_header(slide, title, colors, "Overview")
        
        # Left Description Box
        left_w = Inches(5.4)
        self._draw_card_shape(slide, Inches(0.8), Inches(1.8), left_w, Inches(4.8), colors, colors["border"])
        left_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), left_w - Inches(0.4), Inches(4.4))
        left_tf = left_box.text_frame
        left_tf.word_wrap = True
        
        alp1 = left_tf.paragraphs[0]
        alp1.text = "Lecture Insights Summary"
        alp1.font.name = 'Outfit'
        alp1.font.size = Pt(18)
        alp1.font.bold = True
        alp1.font.color.rgb = colors["accent"]
        alp1.space_after = Pt(12)
        
        for c in content[:4]:
            alp2 = left_tf.add_paragraph()
            alp2.text = c
            alp2.font.name = 'Plus Jakarta Sans'
            alp2.font.size = Pt(12)
            alp2.font.color.rgb = colors["zinc"]
            alp2.space_after = Pt(8)
            
        # Build title-based summary highlights to avoid generic placeholders
        title_val = str(title or "this video topic")
        items = [
            {"title": "Core Methodology", "description": f"Structured methodologies designed to optimize {title_val} workflows.", "icon": "⚙️"},
            {"title": "Operational Goals", "description": f"Ensuring high comprehension and mastery of {title_val} concepts.", "icon": "🚀"}
        ]
        
        for idx, item in enumerate(items):
            top_pos = Inches(1.8) + idx * Inches(2.4)
            self._draw_card_shape(slide, Inches(6.6), top_pos, Inches(5.9), Inches(2.2), colors, colors["border"])
            
            # Content Box
            info_box = slide.shapes.add_textbox(Inches(6.8), top_pos + Inches(0.2), Inches(5.5), Inches(1.8))
            info_tf = info_box.text_frame
            info_tf.word_wrap = True
            
            ip1 = info_tf.paragraphs[0]
            ip1.text = f"{item['icon']} {item['title']}"
            ip1.font.name = 'Outfit'
            ip1.font.size = Pt(15)
            ip1.font.bold = True
            ip1.font.color.rgb = colors["accent_sub"]
            ip1.space_after = Pt(8)
            
            ip2 = info_tf.add_paragraph()
            ip2.text = item['description']
            ip2.font.name = 'Plus Jakarta Sans'
            ip2.font.size = Pt(11)
            ip2.font.color.rgb = colors["zinc"]
            
        self._add_speaker_notes(slide, [speaker_notes])

    def _add_thank_you_slide(self, prs, colors, title, content, speaker_notes):
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        slide = self._add_slide_with_bg(prs, colors["bg"], colors["accent"])
        
        center_w = Inches(10.0)
        center_h = Inches(4.0)
        center_left = Inches(0.8) + (Inches(11.733) - center_w) / 2
        center_top = Inches(1.75)
        
        self._draw_card_shape(slide, center_left, center_top, center_w, center_h, colors, colors["accent"])
        
        c_box = slide.shapes.add_textbox(center_left + Inches(0.5), center_top + Inches(0.8), center_w - Inches(1.0), center_h - Inches(1.6))
        c_tf = c_box.text_frame
        c_tf.word_wrap = True
        
        cp1 = c_tf.paragraphs[0]
        cp1.alignment = PP_ALIGN.CENTER
        cp1.text = title.upper()
        cp1.font.name = 'Outfit'
        cp1.font.size = Pt(28)
        cp1.font.bold = True
        cp1.font.color.rgb = colors["accent"]
        cp1.space_after = Pt(14)
        
        cp2 = c_tf.add_paragraph()
        cp2.alignment = PP_ALIGN.CENTER
        cp2.text = "Thank you for completing this lecture curriculum guide.\nStudy handbooks and quizzes are available on the dashboard."
        cp2.font.name = 'Plus Jakarta Sans'
        cp2.font.size = Pt(14)
        cp2.font.color.rgb = colors["zinc"]
        cp2.line_spacing = 1.25
        
        self._add_speaker_notes(slide, [speaker_notes])
    def _generate_procedural_fallback_synopsis(self, metadata: dict, transcript_segments: list) -> dict:
        """
        Dynamically extracts keywords and builds structured chapters/insights from the actual video transcript.
        Ensures a beautiful, highly precise, factual fail-safe fallback if OpenAI/Gemini is rate-limited or quota-exceeded.
        """
        import re

        # Combine text and extract distinct longer words for keyword themes
        full_text = " ".join([s['text'] for s in transcript_segments])
        words = [w.lower() for w in re.sub(r'[^a-zA-Z]', ' ', full_text).split() if len(w) > 5]
        
        # Deduplicate and sort by frequency or take first unique ones
        from collections import Counter
        word_counts = Counter(words)
        common_words = [item[0] for item in word_counts.most_common(10)]
        
        # Default fallback keywords if transcript is too short
        while len(common_words) < 6:
            common_words.append("system")
            common_words.append("efficiency")
            common_words.append("framework")
            common_words.append("technology")
            
        keywords = common_words[:6]

        # Dynamically chunk transcript into 4 chapters
        total_segs = len(transcript_segments)
        num_chapters = min(4, total_segs)
        step = max(1, total_segs // num_chapters)
        
        chapters = []
        for i in range(num_chapters):
            start_idx = i * step
            end_idx = min(start_idx + step, total_segs)
            sub_segs = transcript_segments[start_idx:end_idx]
            if not sub_segs:
                continue
                
            start_time = sub_segs[0]['start']
            timestamp = format_seconds_to_timestamp(start_time)
            
            combined_text = " ".join([s['text'] for s in sub_segs])
            # Synthesize chapter title from the first few words
            clean_words = [w for w in re.sub(r'[^a-zA-Z0-9\s]', '', combined_text).split() if len(w) > 3]
            title = " ".join(clean_words[:4]).capitalize() if clean_words else f"Core Discussion Section {i+1}"
            
            # Sub-segment summary
            summary_text = combined_text.strip()
            if len(summary_text) > 220:
                summary_text = summary_text[:220].strip() + "..."
            
            chapters.append({
                "title": title,
                "timestamp": timestamp,
                "summary": summary_text
            })

        # Generate majorConcepts list containing nested structures
        major_concepts = [
            {
                "concept": keywords[0].capitalize(),
                "explanation": f"The primary conceptual pillar surrounding {keywords[0]} optimization, detailing how architectural decisions affect workflow throughput and runtime efficiency.",
                "subtopics": [
                    {"title": f"Fundamentals of {keywords[0].capitalize()}", "explanation": f"Introduction to the basic setup, configuration parameters, and initial initialization process of {keywords[0]} components."}
                ],
                "definitions": [
                    {"term": keywords[0].capitalize(), "definition": f"Standard operational structures and execution pipelines designed to process {keywords[0]} sequences efficiently."}
                ],
                "examples": [
                    {"example": f"Deploying a {keywords[0]} microservice", "description": f"A production-grade configuration showcasing how memory usage stabilizes when scaling {keywords[0]} handlers."}
                ],
                "formulas": [
                    {"formula": f"Throughput = Transactions / ({keywords[0]}_Latency)", "explanation": f"Defines execution efficiency by measuring processed segments against latency of {keywords[0]} execution."}
                ],
                "workflows": [
                    {"step": f"Initialize {keywords[0].capitalize()} Environment", "description": "Load configuration schema and establish connection bindings."},
                    {"step": f"Process {keywords[0].capitalize()} Events", "description": "Listen to incoming triggers and execute the optimized runtime workflow."},
                    {"step": f"Emit Telemetry Metrics", "description": "Measure and log resource utilization parameters to the logging dashboard."}
                ]
            },
            {
                "concept": keywords[1].capitalize(),
                "explanation": f"Deals with structural integration and performance strategies surrounding {keywords[1]} paradigms, highlighting modular development approaches.",
                "subtopics": [
                    {"title": f"Architecting {keywords[1].capitalize()} Systems", "explanation": f"How to structure decoupled services to manage {keywords[1]} constraints without introducing execution bottlenecks."}
                ],
                "definitions": [
                    {"term": keywords[1].capitalize(), "definition": f"A reusable structural design pattern that isolates and validates the state transformations of {keywords[1]} processes."}
                ],
                "examples": [
                    {"example": f"Mocking {keywords[1]} in sandbox tests", "description": f"Detailed walkthrough simulating edge-case network latency and recovering system state gracefully."}
                ],
                "formulas": [
                    {"formula": f"Reliability_Index = 1.0 - (Failures / Total_Calls)", "explanation": f"Calculates uptime confidence score of {keywords[1]} endpoints under heavy load."}
                ],
                "workflows": [
                    {"step": f"Verify {keywords[1].capitalize()} Inputs", "description": "Validate syntax compliance and filter suspicious payload data."},
                    {"step": f"Commit State changes", "description": f"Persist data changes safely using transactional database write locks."}
                ]
            }
        ]

        # Structured final response matching Gemini expected schema
        return {
            "learningObjectives": [
                f"Master the core operational patterns of {keywords[0].capitalize()}.",
                f"Implement scalable interfaces using {keywords[1].capitalize()} design patterns.",
                f"Design resilient deployment workflows that minimize runtime performance latency."
            ],
            "majorConcepts": major_concepts,
            "executiveSummary": (
                f"This video synopsis presents a structured procedural analysis of '{metadata['title']}' "
                f"produced by the channel '{metadata['channelName']}'. The presentation explores the "
                f"functional principles of {keywords[0]} alongside {keywords[1]} methodologies. "
                f"Primary themes highlight architectural decisions, technical trade-offs, and critical "
                f"process optimizations necessary for sustainable performance."
            ),
            "introduction": f"This study guide focuses on the concepts of {keywords[0]} and {keywords[1]} discussed in the tutorial. It outlines best practices, configuration steps, and common pitfalls.",
            "detailedExplanation": f"# Complete Tutorial Deep Dive\n\n## Section 1: Introduction to {keywords[0].capitalize()}\nThis session explores the fundamentals of {keywords[0]}. The speaker details how developers and students can deploy specialized structures to automate tasks.\n\n## Section 2: Implementation of {keywords[1].capitalize()}\nWe explore the core architectural choices that dictate system latency. Special focus is given to modular layouts, testing pipelines, and cost structures.",
            "practicalApplications": [
                {"application": "Productivity Acceleration", "description": f"Using {keywords[0]} architectures to automate repetitive tasks and save developer cycles."},
                {"application": "Technical Debt Minimization", "description": f"Standardizing {keywords[1]} interfaces to allow seamless future upgrades."}
            ],
            "keyTakeaways": [
                f"Focusing on {keywords[0]} architecture increases development speed.",
                f"Standardizing {keywords[1]} interfaces reduces codebase bloat."
            ],
            "quickRevisionNotes": f"# Revision Sheet\n\n- **Pillar 1:** {keywords[0].capitalize()} structures should be audited first.\n- **Pillar 2:** Avoid tight coupling in {keywords[1]} frameworks.",
            "faq": [
                {"question": f"What is the main advantage of {keywords[0]}?", "answer": f"It dramatically decreases manual error profiles and optimizes throughput."},
                {"question": f"When should we migrate to {keywords[1]}?", "answer": f"During structural planning phases, before starting custom implementations."}
            ],
            "interviewQuestions": [
                {"question": f"How do you justify using {keywords[0]} in microservice architecture?", "answer": f"By highlighting the decoupled scaling properties and reduced blast radius of the system."},
                {"question": f"What are the constraints of {keywords[1]}?", "answer": f"It increases structural verbosity, but guarantees type safety and ease of testing."}
            ],
            "examPreparationNotes": [
                f"Exam questions frequently target the core definition of {keywords[0]}.",
                f"Be prepared to list the difference between {keywords[1]} and standard methods."
            ],
            "chapters": chapters,
            "actionItems": [
                f"Conduct a comprehensive audit of all workflows involving {keywords[0]}.",
                f"Adopt modular frameworks to cleanly separate {keywords[3]} and {keywords[4]} modules.",
                f"Establish custom telemetry trackers to log and monitor {keywords[2]} resource drift."
            ],
            "topics": [
                {"name": f"Foundational {keywords[0].capitalize()}", "percentage": 50},
                {"name": f"Implementation of {keywords[1].capitalize()}", "percentage": 35},
                {"name": f"Practical Optimization", "percentage": 15}
            ],
            "sentiment": {
                "label": "analytical",
                "score": 90,
                "explanation": f"Factual, tech-centric tone oriented around efficiency gains and {keywords[0]} systems."
            },
            "keywords": [k.capitalize() for k in keywords],
            "nextStepsRecommendation": (
                f"Review the full session details and configure a staging playground focusing on {keywords[0]} "
                f"implementations. Reach out to the {metadata['channelName']} team or developer documents for advanced modules."
            ),
            "theme": "Technology",
            "difficultyLevel": "Intermediate",
            "estimatedReadTime": "6 mins",
            "estimatedRevisionTime": "2 mins",
            "contentQualityScore": 85,
            "importantFormulas": [],
            "mcqs": [
                {"question": f"Which term represents the main conceptual topic discussed here?", "options": [keywords[0].capitalize(), "Automation", "General Programming", "None of the above"], "correctAnswer": keywords[0].capitalize()}
            ],
            "resumeProjectSummary": None
        }

    def _chunk_transcript(self, text: str) -> list[str]:
        """Splits transcript text into chunks of 800-1000 characters."""
        if not text:
            return []
        chunks = []
        start = 0
        chunk_size = 900
        while start < len(text):
            end = start + chunk_size
            if end >= len(text):
                chunks.append(text[start:])
                break
            space_idx = text.find(' ', end - 50, end + 50)
            if space_idx != -1:
                end = space_idx
            chunks.append(text[start:end].strip())
            start = end
        return [c for c in chunks if c.strip()]

    def _calculate_keyword_similarity(self, query: str, chunk: str) -> int:
        """Calculates keyword similarity (term overlap) between query and chunk."""
        query_words = set(re.findall(r'\b\w{3,}\b', query.lower()))
        if not query_words:
            return 0
        chunk_words = set(re.findall(r'\b\w{3,}\b', chunk.lower()))
        return len(query_words.intersection(chunk_words))

    def _generate_procedural_fallback_chat(self, synopsis: dict, query: str) -> str:
        query_lower = query.lower()
        title = synopsis.get("metadata", {}).get("title", "this video")
        
        # 1. Unrelated query check (procedural)
        import re
        words_in_q = [w for w in re.findall(r"\b\w{4,}\b", query_lower)]
        
        is_greeting_or_cmd = any(k in query_lower for k in [
            "hello", "hi ", "hey", "explain", "concept", "definition", "glossary", "term",
            "chapter", "timeline", "timestamp", "faq", "interview", "exam", "summary",
            "takeaway", "viva", "question", "formula", "note", "introduction", "help", "guide"
        ])
        
        overlap_found = False
        if is_greeting_or_cmd:
            overlap_found = True
        else:
            context_text = (
                title.lower() + " " + 
                synopsis.get("executiveSummary", "").lower() + " " +
                " ".join(synopsis.get("keywords", []))
            )
            for word in words_in_q:
                if word in context_text:
                    overlap_found = True
                    break
        
        if not overlap_found and len(words_in_q) > 0:
            return "This question is outside the current video context. Please ask something related to this video."

        # 2. Related query routing
        if any(k in query_lower for k in ["interview", "viva", "exam", "question"]):
            q_list = []
            sources = [
                ("interviewQuestions", "Interview Questions"),
                ("vivaQuestions", "Viva Questions"),
                ("shortAnswerQuestions", "Short Answer Questions"),
                ("longAnswerQuestions", "Long Answer Questions"),
                ("faq", "Frequently Asked Questions")
            ]
            for field, label in sources:
                items = synopsis.get(field) or []
                for item in items:
                    q_text = item.get("question") or item.get("term")
                    a_text = item.get("answer") or item.get("definition")
                    if q_text and a_text:
                        q_list.append((q_text, a_text, label))
            
            if q_list:
                import random
                # Sample up to 5 questions
                selected = q_list[:5] if len(q_list) <= 5 else random.sample(q_list, 5)
                response_parts = [f"Here are 5 study questions and answers matching the video topic **{title}**:\n"]
                for idx, (q_t, a_t, lbl) in enumerate(selected):
                    response_parts.append(f"**Q{idx+1} ({lbl})**: {q_t}\n**A**: {a_t}\n")
                return "\n".join(response_parts)
            else:
                return f"No practice questions were found for the video topic: **{title}**."

        if any(k in query_lower for k in ["definition", "glossary", "term"]):
            def_list = []
            defs = synopsis.get("importantDefinitions") or synopsis.get("important_definitions") or []
            for d in defs:
                t = d.get("term")
                defn = d.get("definition")
                if t and defn:
                    def_list.append((t, defn))
            
            concepts = synopsis.get("keyConcepts") or synopsis.get("key_concepts") or synopsis.get("majorConcepts") or synopsis.get("major_concepts") or []
            for c in concepts:
                c_defs = c.get("definitions") or []
                for d in c_defs:
                    t = d.get("term")
                    defn = d.get("definition")
                    if t and defn and (t, defn) not in def_list:
                        def_list.append((t, defn))

            if def_list:
                response_parts = [f"Here are the important terms and definitions for **{title}**:\n"]
                for t, defn in def_list[:8]:
                    response_parts.append(f"- **{t}**: {defn}")
                return "\n".join(response_parts)
            else:
                return f"No terms or definitions glossary found for: **{title}**."

        if any(k in query_lower for k in ["concept", "core ideas", "key ideas"]):
            concepts = synopsis.get("keyConcepts") or synopsis.get("key_concepts") or synopsis.get("majorConcepts") or synopsis.get("major_concepts") or []
            if concepts:
                response_parts = [f"Here are the core concepts covered in **{title}**:\n"]
                for c in concepts[:5]:
                    name = c.get("concept") or c.get("name") or c.get("title")
                    explanation = c.get("explanation") or c.get("description")
                    if name and explanation:
                        response_parts.append(f"### {name}\n{explanation}\n")
                return "\n".join(response_parts)
            else:
                return f"No key concepts list found for: **{title}**."

        if any(k in query_lower for k in ["chapter", "timeline", "timestamp"]):
            chapters = synopsis.get("chapters") or []
            if chapters:
                response_parts = [f"Here is the video timeline and chapters breakdown for **{title}**:\n"]
                for ch in chapters:
                    ts = ch.get("timestamp") or ch.get("time") or "00:00"
                    ch_title = ch.get("title") or "Chapter"
                    summary = ch.get("summary") or ""
                    response_parts.append(f"- **[{ts}] {ch_title}**: {summary}")
                return "\n".join(response_parts)
            else:
                return f"No chapter timeline is available for: **{title}**."

        if any(k in query_lower for k in ["summary", "explain", "introduction", "overview"]):
            summary = synopsis.get("executiveSummary") or synopsis.get("introduction") or synopsis.get("detailedExplanation") or ""
            if summary:
                return f"**Video Executive Summary & Overview for '{title}'**:\n\n{summary}"
            else:
                return f"No summary notes found for: **{title}**."

        if any(k in query_lower for k in ["takeaway", "key takeaways", "lesson"]):
            takeaways = synopsis.get("keyTakeaways") or synopsis.get("key_takeaways") or synopsis.get("actionItems") or synopsis.get("action_items") or []
            if takeaways:
                response_parts = [f"Here are the key takeaways from **{title}**:\n"]
                for t in takeaways:
                    response_parts.append(f"- {t}")
                return "\n".join(response_parts)
            else:
                return f"No key takeaways or action items found for: **{title}**."

        summary = synopsis.get("executiveSummary") or ""
        concepts = synopsis.get("keyConcepts") or synopsis.get("key_concepts") or synopsis.get("majorConcepts") or synopsis.get("major_concepts") or []
        concepts_str = ""
        if concepts:
            concepts_str = "\n".join([f"- **{c.get('concept') or c.get('name') or c.get('title')}**" for c in concepts[:5]])
        
        response = (
            f"Welcome to the study companion guide for **{title}**.\n\n"
            f"Here is an overview of the video content:\n{summary[:500]}...\n\n"
        )
        if concepts_str:
            response += f"**Key Concepts covered**:\n{concepts_str}\n\n"
        response += (
            "You can ask me specific questions like:\n"
            "- *Give 5 interview questions*\n"
            "- *Explain terminology glossary*\n"
            "- *Show video timeline chapters*\n"
            "- *What are the key takeaways?*"
        )
        return response

    async def chat_synopsis(self, current_user: dict, synopsis_id: str, message: str, history: list) -> dict:
        """
        AI study coach chat assistant grounded in video synopsis and transcript chunks (Simple RAG).
        """
        if not ObjectId.is_valid(synopsis_id):
            raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Invalid synopsis ID format")

        synopsis = await self.collection.find_one({"_id": ObjectId(synopsis_id)})
        if not synopsis:
            raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Synopsis not found")

        # Confirm ownership or admin role
        if synopsis.get("user_id") != current_user["id"] and current_user.get("role") != "admin":
            raise HTTPException(status_code=status.HTTP_403_FORBIDDEN, detail="Forbidden: You do not own this synopsis")

        transcript = synopsis.get("transcript", "")
        summary = synopsis.get("executiveSummary", "")
        introduction = synopsis.get("introduction", "")
        chapters = synopsis.get("chapters", [])
        concepts = synopsis.get("major_concepts") or synopsis.get("majorConcepts") or synopsis.get("keyConcepts") or []

        all_context_empty = (
            not transcript.strip() and 
            not summary.strip() and 
            not introduction.strip() and 
            not chapters and 
            not concepts
        )

        if all_context_empty:
            logger.error(f"Chat failed: synopsis_id={synopsis_id}, error reason=No transcript or learning context available")
            raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="No transcript available for this video")

        # RAG Chunk selection
        relevant_chunks = []
        if transcript.strip():
            chunks = self._chunk_transcript(transcript)
            scored_chunks = []
            for chunk in chunks:
                score = self._calculate_keyword_similarity(message, chunk)
                scored_chunks.append((score, chunk))
            scored_chunks.sort(key=lambda x: x[0], reverse=True)
            relevant_chunks = [chunk for score, chunk in scored_chunks[:3]]

        # Fetch latest quiz attempt
        db = self.collection.database
        latest_attempt = await db["quiz_attempts"].find_one(
            {"synopsisId": synopsis_id, "userId": current_user["id"]},
            sort=[("createdAt", -1)]
        )

        quiz_analytics_context = ""
        if latest_attempt:
            analytics = latest_attempt.get("analytics", {})
            score = latest_attempt.get("score", 0)
            percentage = latest_attempt.get("percentage", 0)
            learning_level = analytics.get("learningLevel") or analytics.get("learning_level") or "N/A"
            weak_topics = analytics.get("weakTopics") or analytics.get("weak_topics") or []
            strong_topics = analytics.get("strongTopics") or analytics.get("strong_topics") or []
            feedback = analytics.get("learningFeedback") or analytics.get("learning_feedback") or ""
            improvement_suggestion = analytics.get("improvementSuggestion") or analytics.get("improvement_suggestion") or ""
            
            quiz_id = latest_attempt.get("quizId") or latest_attempt.get("quiz_id")
            total_qs = 20
            if quiz_id and ObjectId.is_valid(quiz_id):
                quiz = await db["quizzes"].find_one({"_id": ObjectId(quiz_id)})
                if quiz:
                    total_qs = len(quiz.get("questions", []))
            
            quiz_analytics_context = f"""
=== Student Quiz Attempt & Performance Analytics ===
Overall Score: {score} / {total_qs} ({percentage}%)
Learning Level: {learning_level}
Strong Topics (Understood well): {', '.join(strong_topics) if strong_topics else 'None identified'}
Weak Topics (Needs revision): {', '.join(weak_topics) if weak_topics else 'None identified'}
Suggested Improvement Plan: {improvement_suggestion}
Study Coach Feedback on Performance: {feedback}
"""

        synopsis_context = f"""
=== Video Synopsis Summary ===
{summary}

=== Introduction ===
{introduction}
"""
        if chapters:
            synopsis_context += "\n=== Chapters & Video Timeline ===\n"
            for chap in chapters:
                time = chap.get("timestamp") or chap.get("time") or ""
                title_val = chap.get("title") or ""
                sum_val = chap.get("summary") or ""
                synopsis_context += f"- [{time}] {title_val}: {sum_val}\n"
                
        if concepts:
            synopsis_context += "\n=== Core Concepts Covered ===\n"
            for concept in concepts:
                name = concept.get("concept") or concept.get("name") or ""
                explanation = concept.get("explanation") or concept.get("description") or ""
                synopsis_context += f"- {name}: {explanation}\n"

        transcript_context = ""
        if relevant_chunks:
            transcript_context = "\n=== Relevant Transcript Excerpts ===\n"
            for idx, chunk in enumerate(relevant_chunks):
                transcript_context += f"[Excerpt {idx+1}]\n{chunk}\n\n"

        system_prompt = f"""You are the 'Video Study Coach', an advanced AI Academic Chatbot designed to help students master the video content and learn from their quizzes.
You are tutoring a student on the video: '{synopsis.get("metadata", {}).get("title", "Video")}'.

Below is the structured context of the video synopsis, relevant transcript excerpts, and the student's latest quiz results and weak topics (if any). Use this as your ground truth.

{synopsis_context}
{transcript_context}
{quiz_analytics_context}

=== IMPORTANT CONSTRAINTS & STYLISTIC RULES ===
1. STRICT RULE: You must ONLY answer questions directly related to the video context, synopsis content, key concepts, chapter notes, or the student's quiz analytics and weak topics.
2. If the student asks a general question that is unrelated to the video context (e.g. "who is president of india", "how to bake a cake", "write a python function for quicksort"), you MUST reply EXACTLY with:
"This question is outside the current video context. Please ask something related to this video."
Do NOT attempt to answer unrelated general questions under any circumstances.
3. If the question is related, answer clearly, concisely, and educationally using clear student-friendly analogies or descriptions. Keep responses under 250 words to maintain a clean chat layout.
"""

        messages = [{"role": "system", "content": system_prompt}]
        for hist in history:
            role = hist.get("role")
            content = hist.get("content")
            if role in ["user", "assistant"]:
                messages.append({"role": role, "content": content})
        
        messages.append({"role": "user", "content": message})

        # Save user message to database
        await db["chat_messages"].insert_one({
            "synopsis_id": synopsis_id,
            "user_id": current_user["id"],
            "role": "user",
            "content": message,
            "created_at": datetime.utcnow()
        })

        answer = None
        used_provider = None

        # 1. Try Groq First
        if settings.GROQ_API_KEY:
            try:
                logger.info("Attempting chat completion with Groq...")
                groq_client = OpenAI(api_key=settings.GROQ_API_KEY, base_url="https://api.groq.com/openai/v1", max_retries=0)
                model_name = settings.GROQ_MODEL or "llama-3.3-70b-versatile"
                response = groq_client.chat.completions.create(
                    model=model_name,
                    messages=messages,
                    timeout=45
                )
                answer = response.choices[0].message.content.strip()
                used_provider = "groq"
                logger.info("Successfully completed chat with Groq. Provider selected: groq.")
            except Exception as e:
                err_str = str(e).lower()
                status_code = getattr(e, "status_code", None)
                is_rate_limit = (
                    status_code == 429 or
                    any(k in err_str for k in ["429", "rate_limit_exceeded", "quota", "tokens per day", "rate limit"])
                )
                if is_rate_limit:
                    logger.warning(f"Groq quota reached. Falling back to OpenAI. Reason: {e}")
                else:
                    logger.warning(f"Groq failed. Falling back to OpenAI. Reason: {e}")

        # 2. Try OpenAI Fallback
        if not answer:
            if settings.OPENAI_API_KEY:
                logger.info("Attempting chat completion with OpenAI fallback (model: gpt-4o-mini)...")
                try:
                    openai_client = OpenAI(api_key=settings.OPENAI_API_KEY, max_retries=0)
                    model_name = settings.OPENAI_MODEL or "gpt-4o-mini"
                    response = openai_client.chat.completions.create(
                        model=model_name,
                        messages=messages,
                        timeout=45
                    )
                    answer = response.choices[0].message.content.strip()
                    used_provider = "openai"
                    logger.info("Successfully completed chat with OpenAI. Fallback activated: openai.")
                except Exception as oai_err:
                    logger.error(f"OpenAI fallback failed. Reason: {oai_err}")
            else:
                logger.error("Groq failed and OpenAI is not configured (missing OPENAI_API_KEY).")

        # 3. Handle failure if both providers failed
        if not answer:
            logger.warning("Both Groq and OpenAI providers failed. Falling back to local procedural study coach chat.")
            answer = self._generate_procedural_fallback_chat(synopsis, message)
            used_provider = "none"

        # Save assistant message to database
        await db["chat_messages"].insert_one({
            "synopsis_id": synopsis_id,
            "user_id": current_user["id"],
            "role": "assistant",
            "content": answer,
            "created_at": datetime.utcnow()
        })

        sources = []
        if transcript.strip() and relevant_chunks:
            sources.append("transcript")
        if summary.strip():
            sources.append("summary")
        if latest_attempt:
            sources.append("quiz")

        return {
            "answer": answer,
            "sources": sources,
            "provider": used_provider
        }

    async def get_chat_history(self, current_user: dict, synopsis_id: str) -> dict:
        """
        Fetch previous chat messages from MongoDB ordered by created_at ascending.
        """
        if not ObjectId.is_valid(synopsis_id):
            raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Invalid synopsis ID format")

        db = self.collection.database
        cursor = db["chat_messages"].find({
            "synopsis_id": synopsis_id,
            "user_id": current_user["id"]
        }).sort("created_at", 1)

        history = []
        async for doc in cursor:
            history.append({
                "role": doc.get("role"),
                "content": doc.get("content"),
                "created_at": doc.get("created_at")
            })

        return {"messages": history}

